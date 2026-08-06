#!/usr/bin/env python3
"""
Create a 10-slide PowerPoint presentation on complex topics with speaker notes.
Organized as Sourov's life narrative arc.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """Create and save the presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define slides: (title, bullet_points, speaker_notes)
    slides_data = [
        {
            "title": "Seven Concepts, One Life:\nThriving in Uncertainty",
            "bullets": [
                "Psychology: Productivity after trauma & ADHD",
                "Linguistics: Phonology shapes identity & memory",
                "Physics: Quantum thinking meets meditation",
                "AI: What LLM limitations teach us about human learning",
                "Personal: How these concepts interconnect in one journey"
            ],
            "notes": """You'll explore seven rare-angle concepts through one person's real life journey—from Bangladesh to France to Australia, through workplace trauma to recovery, from language learning to AI. The promise: abstract ideas become concrete when they illuminate lived experience. This isn't a textbook. It's a story that teaches."""
        },
        {
            "title": "From Bangladesh to France:\nWhy We Belong (and Where)",
            "bullets": [
                "Evolutionary bias: unconscious patterns in migration choice",
                "Not destiny, but recognition: naming the pattern enables choice",
                "Sourov's journey: Why leave Bangladesh? Why France? Why Australia?",
                "Agency over instinct: conscious values override evolutionary pushes",
                "The insight: Understanding bias is step 1 to choosing differently"
            ],
            "notes": """Migration isn't purely rational calculation—there's an evolutionary substrate pulling us toward status, security, belonging. Sourov's move from Bangladesh to France wasn't accidental; recognizing the patterns (security-seeking, cultural capital-chasing) helped him later choose more consciously about Australia. The rare angle: evolutionary bias isn't something to shame yourself for; it's something to see clearly, then transcend."""
        },
        {
            "title": "Code-Switching:\nHow Languages Shape Who We Are",
            "bullets": [
                "Definition: Phonological + grammatical shifts based on social context",
                "Not 'fakeness': It's how all multilingual people navigate belonging",
                "Sourov's experience: Bangladesh accent → French school → Australia moves accent again",
                "Identity emerges in how you speak, not despite it",
                "The cost and the freedom: Code-switching as both labor and power"
            ],
            "notes": """Code-switching isn't a failure to integrate or an act of deception. Linguists study it as evidence of sophisticated cognitive flexibility—you're managing multiple identity-contexts simultaneously. Sourov's journey across three countries shows code-switching in real time: accent, vocabulary, humor all shift. The rare angle: phonological change is identity work, and it's exhausting and powerful at once."""
        },
        {
            "title": "The Sound of Memory:\nWhy Phonology Sticks",
            "bullets": [
                "Phonological priming: Certain sounds activate memory without awareness",
                "Universal patterns: Phonology works across language boundaries",
                "Why this matters: You use phonological memory every day (brands, songs, names)",
                "Example: Why 'Coca-Cola' sticks even in languages without those sounds",
                "The insight: Cognition is shaped by sound patterns you never consciously notice"
            ],
            "notes": """Phonological patterns don't just help us pronounce—they shape how memory itself works. Certain sound sequences feel 'right' across cultures because they exploit universal cognitive constraints. Sourov's teaching career shows this: pronunciation drills are actually memory drills. The rare angle: phonology is a cognitive lever everyone pulls, but linguists study it as though it's exotic. It's not—it's how your brain works."""
        },
        {
            "title": "Productivity Lost & Found:\nADHD, Trauma, and Recovery",
            "bullets": [
                "The breakdown: Workplace conflict + trauma → loss of productivity confidence",
                "ADHD reframing: Not a deficit, but a toolkit for post-trauma recovery",
                "Three protocols: Clear externals (lists, timers, environment), flexible internally (permission to rest, non-linear progress), pattern-seeking as asset (ADHD brains see connections)",
                "Trauma-informed insight: Recovery isn't 'getting back to normal'; it's building a new system",
                "Sourov's discovery: ADHD traits that felt like liabilities became the recovery toolkit"
            ],
            "notes": """Post-trauma productivity loss is real and isn't fixed by 'hustle harder.' But ADHD minds often have unexpected resilience post-trauma because they've built external scaffolding (lists, timers, accountability) for their attention. Sourov's workplace trauma was compounded by undiagnosed ADHD; recovery came when he stopped fighting the ADHD and started using it. The rare angle: ADHD as pragmatic asset, not inspirational narrative."""
        },
        {
            "title": "Systems Thinking for Uncertain Minds:\nClear Rules, Flexible Tactics",
            "bullets": [
                "System prompt concept (AI): Consistent instructions + adaptive execution = stability",
                "Apply to humans: Your values are your 'system prompt'; your daily tactics are adaptive",
                "For ADHD & trauma-recovered minds: External structure + internal flexibility = freedom",
                "Example: Sourov's stability routine—same time to start, permission to adjust during the day",
                "The meta-insight: Stability comes from clarity about what doesn't change, not control of everything"
            ],
            "notes": """LLMs are stable because they have consistent parameters (system prompt, temperature, top-k) combined with flexible outputs. Human stability works the same way—you need crystal-clear values/rules, then permission to execute flexibly. This matters for ADHD brains especially, which struggle with rigid approaches but thrive with clear principles + tactical freedom. The teaching: write your system prompt for your own life."""
        },
        {
            "title": "Teaching Machines What Humans Know:\nLLMs as Language Learners",
            "bullets": [
                "LLMs learn statistics of language patterns; humans learn language embodied (taste, touch, social stakes)",
                "What LLMs can't do: Connect words to lived experience, understand cultural context, know when to break rules",
                "What this teaches us: Human language learning is deeper than pattern recognition—it's cultural apprenticeship",
                "Sourov's ELT insight: Teaching pronunciation isn't teaching sounds; it's teaching belonging",
                "The rare angle: LLM limitations illuminate what human language teachers actually do"
            ],
            "notes": """People often worry 'AI will replace language teachers.' But watching LLMs struggle with nuance and context makes clear what human teachers actually teach: not pattern-matching, but judgment, culture, identity, belonging. Sourov's students aren't learning English; they're learning how to be heard in English. LLMs can't do that. The teaching: understanding AI limitations is the best way to understand human expertise."""
        },
        {
            "title": "Superposition to Meditation:\nEmbracing Uncertainty",
            "bullets": [
                "Quantum superposition: Two states at once; measured → collapses to one",
                "Metaphor for ADHD attention: Hyperaware + scattered simultaneously (not sequential)",
                "Meditation practice: Observing attention without collapsing it into single focus",
                "The wisdom: Certainty is a trap; probability (superposition) is how reality works",
                "Sourov's insight: ADHD wasn't a bug until he started meditating to 'fix' it; acceptance changed everything"
            ],
            "notes": """This is a metaphor, not physics—but a powerful one. ADHD attention IS like superposition: you're genuinely attending to multiple threads simultaneously, and the moment you 'force focus' you lose the multi-threaded insight. Meditation teaches you to observe attention without collapsing it. Sourov's meditation practice made peace with ADHD not by 'fixing' it but by accepting how his attention actually works. The rare angle: the problem isn't ADHD; it's fighting ADHD."""
        },
        {
            "title": "Evolutionary Bias in Love, Marriage, Belonging",
            "bullets": [
                "Evolutionary psychology: Partnership choices follow unconscious patterns (status, health markers, family structure)",
                "Sourov's marriage: Conscious choice about values despite evolutionary pulls",
                "Migration to Australia: Same pattern—evolutionary drive for security meets conscious vision for life",
                "The insight: Recognizing bias doesn't eliminate it; it enables choosing despite it",
                "Agency lives in the gap between instinct and consciousness"
            ],
            "notes": """We often tell ourselves partnerships and migrations are purely rational decisions. They're not. But they're also not purely instinctive. Evolutionary psychology maps the unconscious terrain; consciousness is what lets you choose your path through it. Sourov's journey—why marry this person, why move to Australia—shows both: deep patterns recognized, then transcended by conscious choice. The teaching: you're not a puppet of evolution, but understanding the puppet-strings is the first step to freedom."""
        },
        {
            "title": "Building Your Stability System:\n3 Protocols",
            "bullets": [
                "Protocol 1: Know your system prompt (your non-negotiable values)",
                "Protocol 2: Build external scaffolding (lists, routines, environment design for your neurology)",
                "Protocol 3: Practice observation without judgment (meditation, self-awareness, recognizing patterns without fighting them)",
                "Integration: How to recognize rare angles in your own story and apply them",
                "Next step: Write your 'system prompt'; test your external scaffolding; practice 10 minutes daily"
            ],
            "notes": """The entire presentation isn't about Sourov—it's about how to build stability in your own life using these concepts. Take the ADHD framework, the phonological awareness, the evolutionary insight, the quantum acceptance—and make them yours. A system prompt for your life is crystal clarity about what matters, combined with permission to execute tactically. Your external scaffolding is whatever works for your neurology (lists, routines, space, time-blocking). Your observation practice is recognizing patterns without fighting them. That combination = stability."""
        }
    ]

    # Add title slide (layout 6 is typically blank)
    for i, slide_data in enumerate(slides_data):
        # Use blank layout for flexibility
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_p = title_frame.paragraphs[0]
        title_p.text = slide_data["title"]
        title_p.font.size = Pt(44)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(25, 25, 112)  # Midnight blue

        # Add bullet points
        bullets_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.25), Inches(8.5), Inches(4.5))
        bullets_frame = bullets_box.text_frame
        bullets_frame.word_wrap = True

        for idx, bullet in enumerate(slide_data["bullets"]):
            if idx == 0:
                p = bullets_frame.paragraphs[0]
            else:
                p = bullets_frame.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(24)
            p.space_before = Pt(12)
            p.space_after = Pt(6)

        # Add speaker notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = slide_data["notes"]

    return prs

if __name__ == "__main__":
    prs = create_presentation()
    output_path = "/home/user/my_professional_documents/presentations/Sourov_7_Concepts_Psychology_Linguistics_Physics_AI.pptx"
    prs.save(output_path)
    print(f"✓ Presentation created: {output_path}")
