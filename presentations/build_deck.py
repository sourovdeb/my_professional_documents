#!/usr/bin/env python3
"""Builds the 10-slide psychology/linguistics/physics/AI presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette --------------------------------------------------------------
INK = RGBColor(0x1B, 0x1F, 0x2A)          # near-black
PAPER = RGBColor(0xFA, 0xF9, 0xF6)        # warm off-white
ACCENT = RGBColor(0x2E, 0x5E, 0xAA)       # deep blue
ACCENT2 = RGBColor(0xC2, 0x5B, 0x2E)      # burnt orange
MUTED = RGBColor(0x6B, 0x71, 0x80)        # slate gray
LINE = RGBColor(0xDD, 0xD8, 0xCE)

DOMAIN_COLORS = {
    "PSYCHOLOGY": RGBColor(0x2E, 0x5E, 0xAA),
    "LINGUISTICS": RGBColor(0xC2, 0x5B, 0x2E),
    "EVOLUTIONARY BIOLOGY": RGBColor(0x3E, 0x7C, 0x59),
    "ARTIFICIAL INTELLIGENCE": RGBColor(0x7A, 0x4C, 0x9E),
    "QUANTUM PHYSICS": RGBColor(0x1E, 0x7A, 0x8C),
    "SYNTHESIS": RGBColor(0x8C, 0x3A, 0x5B),
}

SLIDES = [
    {
        "kind": "title",
        "eyebrow": "A ten-slide field guide",
        "title": "Minds, Words, Particles, Machines",
        "subtitle": "Rarely-taught ideas from psychology, linguistics, evolutionary biology, AI, and quantum physics",
        "notes": (
            "Welcome the audience and frame the deck as a tour across five fields that all "
            "study some version of the same question: how a system — a brain, a language, a "
            "species, a neural network, or the universe itself — deals with uncertainty and makes "
            "predictions. Mention that every slide is chosen for being genuinely useful or "
            "surprising, not textbook trivia, and that the last slide ties all four fields together."
        ),
    },
    {
        "kind": "content",
        "domain": "PSYCHOLOGY",
        "title": "Why Willpower Isn't the Real Lever",
        "bullets": [
            "Ultradian rhythm: focus naturally cycles in ~90-minute waves (Kleitman); working against the trough, not the peak, is what makes tasks feel impossible.",
            "The Zeigarnik effect: unfinished tasks intrude on memory far more than finished ones — it's why an open loop nags you at 2 a.m.",
            "Implementation intentions (Gollwitzer): a concrete \"if X, then I will Y\" plan roughly doubles follow-through versus a vague goal.",
            "Attention residue (Sophie Leroy): switching tasks before finishing leaves a cognitive trace of the old task, quietly taxing performance on the new one.",
            "Takeaway: productivity systems that ignore biological rhythm and unclosed loops are fighting the wrong battle.",
        ],
        "notes": (
            "Start with the surprising claim that willpower is a much smaller factor in productivity than structure and timing. "
            "Walk through ultradian rhythms and the Zeigarnik effect as evidence that the brain has built-in cycles and unfinished-business "
            "alarms that better protocols work with instead of against. Close by connecting implementation intentions and attention residue "
            "to a practical rule: batch similar work, close loops explicitly, and match hard tasks to your natural peak, not your calendar."
        ),
    },
    {
        "kind": "content",
        "domain": "PSYCHOLOGY",
        "title": "Ego Depletion: A Famous Theory That Mostly Broke",
        "bullets": [
            "Baumeister's 1998 \"ego depletion\" claimed willpower is a finite muscle that tires with use — it became one of psychology's most cited ideas.",
            "A 2016 multi-lab Registered Replication Report (23 labs, 2,000+ participants) found the effect near zero.",
            "The \"glucose restores willpower\" claim also failed to replicate rigorously — the brain doesn't run low on fuel that fast.",
            "What does hold up: motivation, belief about willpower, and switching costs (attention residue) explain most of what looked like \"depletion.\"",
            "Lesson for the audience: a famous, intuitive finding can dominate pop psychology for two decades before failing replication.",
        ],
        "notes": (
            "Use this as a mini case study in the replication crisis — a theory nearly everyone has heard of turned out to be much weaker than believed. "
            "Explain the difference between an original single-lab finding and a large pre-registered multi-lab replication, and why the latter is far more trustworthy. "
            "The practical point for the audience is to be skeptical of tidy metaphors like 'willpower is a muscle' and to look at what actually replicates: motivation and belief matter more than a depleting resource."
        ),
    },
    {
        "kind": "content",
        "domain": "LINGUISTICS",
        "title": "Why English Spelling Feels Broken",
        "bullets": [
            "Coarticulation: neighboring sounds physically bleed into each other, so speech is really a continuous smear, not discrete letter-sounds.",
            "The Great Vowel Shift (c. 1400–1700): English vowel pronunciation changed dramatically after spelling had already been fixed by early printing.",
            "Result: \"knight,\" \"through,\" and \"though\" preserve pronunciations from centuries ago — the spelling is a fossil record, not a bad design.",
            "Infants are universal phoneme listeners at 6 months and can hear contrasts from any language; by ~12 months that ability narrows to their native language's sound categories.",
            "Voice onset time (the tiny delay between a consonant and voicing) is what your brain uses to sort \"b\" from \"p\" — a difference of milliseconds.",
        ],
        "notes": (
            "Open by asking why English spelling seems so irrational, then reveal it's not irrational — it's historical. "
            "Explain the Great Vowel Shift as the moment pronunciation raced ahead of spelling, which had just been locked in by the printing press, leaving today's mismatches as a kind of linguistic fossil. "
            "Then pivot to phonology itself: coarticulation shows speech has no natural letter boundaries, and the infant phoneme-narrowing finding is a striking, lesser-known result showing how early the brain specializes for its native language."
        ),
    },
    {
        "kind": "content",
        "domain": "LINGUISTICS",
        "title": "Does Your Language Bend How You Think?",
        "bullets": [
            "Guugu Yimithirr (Australia) has no words for \"left\" or \"right\" — speakers use absolute directions (north/south/east/west) even for a crumb on a shirt.",
            "Speakers of such languages maintain an accurate sense of cardinal direction at all times, a skill most \"left/right\" language speakers never develop.",
            "Kuuk Thaayorre speakers, when asked to arrange time-sequence cards, order them east-to-west — literally by compass bearing, not by dominant hand.",
            "Russian has separate basic words for light blue (goluboy) and dark blue (siniy); Russian speakers distinguish the two shades faster than English speakers in lab tests.",
            "This is weak Sapir-Whorf: language doesn't cage thought, but it does quietly tune default attention and habitual judgment.",
        ],
        "notes": (
            "Introduce the linguistic relativity debate carefully: the strong version (language imprisons thought) is discredited, but a well-evidenced weak version survives. "
            "The Guugu Yimithirr and Kuuk Thaayorre examples are vivid, well-replicated findings that are rarely taught outside linguistics courses — invite the audience to imagine always knowing compass direction the way they know left and right. "
            "Close with the Russian blues study as a clean, measurable example of the same effect showing up in a simple reaction-time task, not just anecdote."
        ),
    },
    {
        "kind": "content",
        "domain": "EVOLUTIONARY BIOLOGY",
        "title": "Your Biases Are Bugs That Used to Be Features",
        "bullets": [
            "Loss aversion (Kahneman & Tversky): losing $100 hurts roughly twice as much as gaining $100 feels good — an asymmetry that once kept our ancestors from risking scarce resources.",
            "Negativity bias: threats and criticism register faster and stronger in the brain than praise or good news — one angry face pops out of a crowd of happy ones almost instantly.",
            "The smoke-detector principle (Randolph Nesse): evolution tunes alarm systems (fear, nausea, jealousy) to over-fire, because a false alarm is cheap but a missed real threat can be fatal.",
            "In-group bias formed in small bands where trusting strangers was genuinely risky — the same wiring now misfires at the scale of nations and social media.",
            "None of these are irrational in their original context — they're mismatches between ancestral environments and modern ones.",
        ],
        "notes": (
            "Frame cognitive biases not as design flaws but as an evolutionary trade-off: mechanisms optimized for reproductive success in small ancestral groups, not for accuracy in a modern world. "
            "The smoke-detector principle is the key unifying idea — explain the logic of asymmetric error costs, where a system that cries wolf occasionally still wins if it never misses a real wolf. "
            "End by noting this reframes self-improvement: you can't delete these instincts, but understanding their old job helps you catch when they're misapplied today."
        ),
    },
    {
        "kind": "content",
        "domain": "ARTIFICIAL INTELLIGENCE",
        "title": "What an LLM Actually Does, Step by Step",
        "bullets": [
            "Tokenization: text is chopped into sub-word pieces (byte-pair encoding) — \"unbelievable\" might become \"un\" + \"believ\" + \"able,\" not whole words.",
            "Self-attention: for every token, the model computes how much every other token in the context should influence it — this is what lets it track \"it\" back to the right noun across a paragraph.",
            "Training objective is deceptively simple: predict the next token, over and over, across trillions of tokens of text.",
            "Scaling laws: loss (prediction error) falls in a smooth, predictable curve as you add data, parameters, and compute — one of the few genuinely lawlike patterns in deep learning.",
            "Emergent abilities: skills like multi-step arithmetic or chain-of-thought reasoning appear suddenly at certain scale thresholds rather than growing gradually.",
        ],
        "notes": (
            "Ground the audience in the mechanics before the philosophy: tokenization, attention, and next-token prediction are the three ideas that explain almost everything an LLM does. "
            "Emphasize how strange it is that something as simple as 'predict the next token' produces translation, coding, and reasoning behavior once scaled far enough — that gap between the simple objective and the complex behavior is the central mystery of the field. "
            "Mention scaling laws as one of deep learning's few genuinely predictable, almost physics-like regularities, which is part of why labs kept investing in ever-larger models."
        ),
    },
    {
        "kind": "content",
        "domain": "ARTIFICIAL INTELLIGENCE",
        "title": "Hallucination and Grokking: Two Under-Taught LLM Quirks",
        "bullets": [
            "Hallucination isn't a random glitch — it's the same next-token-prediction machinery producing fluent text with no built-in \"I don't know\" signal.",
            "RLHF (human feedback tuning) often rewards confident-sounding answers over honest uncertainty, which can make hallucination worse, not better.",
            "Grokking: some small models trained on simple tasks memorize the training data first, plateau for a long time, then suddenly \"get it\" and generalize perfectly — long after training loss looked flat.",
            "Grokking suggests models can be quietly reorganizing their internal representations during what looks like a boring plateau, not just idling.",
            "In-context learning: a model can learn a brand-new pattern from a few examples in the prompt alone, with no weight updates at all — a capability nobody explicitly trained into it.",
        ],
        "notes": (
            "Hallucination is widely discussed but usually explained badly — clarify that it's a natural consequence of a model trained purely to produce plausible next tokens, not evidence of the model 'lying' or malfunctioning in a special way. "
            "Grokking is a genuinely obscure and delightful research finding worth spending real time on: a model can appear stuck memorizing for thousands of training steps and then abruptly generalize, which challenges the assumption that a flat loss curve means nothing is happening. "
            "In-context learning closes the slide by pointing out that this capability was an unplanned emergent behavior, discovered after the fact, not a designed feature."
        ),
    },
    {
        "kind": "content",
        "domain": "QUANTUM PHYSICS",
        "title": "The Weirdness Is the Rule, Not the Exception",
        "bullets": [
            "Superposition: a quantum particle's state is a combination of possibilities until measured — not \"unknown to us\" but genuinely undetermined.",
            "The double-slit experiment: individual electrons fired one at a time still build up an interference pattern, as though each one passes through both slits and interferes with itself.",
            "Entanglement: two particles can share a linked state such that measuring one instantly determines the other's outcome, regardless of distance — Einstein's \"spooky action at a distance,\" now experimentally confirmed.",
            "Quantum tunneling: particles can pass through energy barriers that classical physics says are impossible to cross — it's the reason the sun can fuse hydrogen at its actual core temperature.",
            "The quantum Zeno effect: frequent enough measurement can literally freeze a quantum system's evolution — a real, lab-verified version of \"a watched pot never boils.\"",
        ],
        "notes": (
            "Set expectations that quantum mechanics isn't just small and hard to see, it behaves in ways that violate everyday intuition about objects having definite properties before you look at them. "
            "Walk through the double-slit result as the clearest demonstration of superposition, then use entanglement to show correlations that have no classical explanation, stressing this has been tested and confirmed, not merely theorized. "
            "The quantum Zeno effect is the rarely-known crowd-pleaser here — spend a moment on it since 'a watched pot never boils' being literally true in a lab is a memorable, discussion-sparking fact."
        ),
    },
    {
        "kind": "content",
        "domain": "SYNTHESIS",
        "title": "One Thread Through All Four Fields",
        "bullets": [
            "The predictive brain: neuroscience increasingly models the brain as a prediction engine, constantly guessing sensory input and updating on error — not a passive recorder of reality.",
            "Language evolved as a low-bandwidth code for compressing and transmitting those predictions between brains, which is why ambiguity, idiom, and metaphor are features, not bugs.",
            "Evolutionary biases are simply predictions that got hard-coded because they paid off often enough, even though they're miscalibrated for modern life.",
            "LLMs formalize the same principle in silicon: a system trained on nothing but \"predict what comes next\" turns out to approximate reasoning, dialogue, and knowledge retrieval.",
            "Even physics bends toward this frame — quantum mechanics is fundamentally about probabilities and measurement, not certainties, all the way down.",
        ],
        "notes": (
            "Bring the whole talk together by naming the shared thread explicitly: prediction under uncertainty shows up as the brain's core algorithm, as the reason language and biases exist, as the entire training objective of an LLM, and as the mathematical structure of quantum mechanics itself. "
            "Encourage the audience to notice this pattern the next time they meet a new field — asking 'what is this system predicting, and what does it do when the prediction is wrong?' is a surprisingly portable lens. "
            "Close by inviting questions and pointing out that each of the nine ideas covered has a rich, well-documented literature if anyone wants to go deeper on a specific slide."
        ),
    },
]

FOOTER_TEXT = "Minds, Words, Particles, Machines"


def set_background(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf


def style_run(run, size, color, bold=False, italic=False, font="Calibri"):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def build_title_slide(slide, data):
    set_background(slide, INK)

    # thin accent rule
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.35), Inches(1.4), Pt(4))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT2
    rule.line.fill.background()

    _, tf = add_textbox(slide, Inches(0.9), Inches(1.75), Inches(11.5), Inches(0.5))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = data["eyebrow"].upper()
    style_run(r, 16, RGBColor(0xC9, 0xCE, 0xD8), bold=True)
    p.alignment = PP_ALIGN.LEFT

    _, tf = add_textbox(slide, Inches(0.85), Inches(2.55), Inches(11.6), Inches(2.0))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = data["title"]
    style_run(r, 48, PAPER, bold=True)
    p.alignment = PP_ALIGN.LEFT

    _, tf = add_textbox(slide, Inches(0.9), Inches(4.15), Inches(10.8), Inches(1.3))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = data["subtitle"]
    style_run(r, 20, RGBColor(0xB8, 0xBF, 0xCC))
    p.alignment = PP_ALIGN.LEFT

    _, tf = add_textbox(slide, Inches(0.9), Inches(6.85), Inches(10.8), Inches(0.4))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Psychology  ·  Linguistics  ·  Evolutionary Biology  ·  Artificial Intelligence  ·  Quantum Physics"
    style_run(r, 13, MUTED, italic=True)
    p.alignment = PP_ALIGN.LEFT


def build_content_slide(slide, data, index, total):
    set_background(slide, PAPER)
    color = DOMAIN_COLORS.get(data["domain"], ACCENT)

    # left accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    # eyebrow: domain + slide number
    _, tf = add_textbox(slide, Inches(0.65), Inches(0.35), Inches(10.5), Inches(0.4))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"{data['domain']}"
    style_run(r, 14, color, bold=True)
    r2 = p.add_run()
    r2.text = f"   ·   {index} / {total}"
    style_run(r2, 14, MUTED)

    # title
    _, tf = add_textbox(slide, Inches(0.6), Inches(0.75), Inches(11.8), Inches(1.1))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = data["title"]
    style_run(r, 30, INK, bold=True)

    # divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.75), Inches(11.4), Pt(1.5))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()

    # bullets
    _, tf = add_textbox(slide, Inches(0.65), Inches(2.05), Inches(11.7), Inches(5.0))
    tf.word_wrap = True
    for i, bullet in enumerate(data["bullets"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(16)
        p.line_spacing = 1.15
        marker = p.add_run()
        marker.text = "―  "
        style_run(marker, 17, color, bold=True)
        r = p.add_run()
        r.text = bullet
        style_run(r, 17, INK)

    # footer
    _, tf = add_textbox(slide, Inches(0.65), Inches(7.05), Inches(11.7), Inches(0.35))
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = FOOTER_TEXT
    style_run(r, 10, MUTED, italic=True)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    content_slides = [s for s in SLIDES if s["kind"] == "content"]

    for data in SLIDES:
        slide = prs.slides.add_slide(blank)
        if data["kind"] == "title":
            build_title_slide(slide, data)
        else:
            idx = content_slides.index(data) + 1
            build_content_slide(slide, data, idx, len(content_slides))

        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = data["notes"]

    out_path = "/home/user/my_professional_documents/presentations/minds_words_particles_machines.pptx"
    prs.save(out_path)
    print("Saved:", out_path)


if __name__ == "__main__":
    main()
