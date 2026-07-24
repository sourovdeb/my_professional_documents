#!/usr/bin/env python3
"""Build 02_Model.pptx for the "AI Explained Simply" series (Episode 2: Model).
Self-contained: only depends on python-pptx. Run: python3 build_deck.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette (Ocean Gradient, matches Episode 1) ----
NAVY = RGBColor(0x21, 0x29, 0x5C)
DEEP_BLUE = RGBColor(0x06, 0x5A, 0x82)
TEAL = RGBColor(0x1C, 0x72, 0x93)
ICE = RGBColor(0xEA, 0xF2, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x16, 0x20, 0x2A)
MUTED = RGBColor(0x5B, 0x6B, 0x76)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame


def para(tf, text, size, color, bold=False, first=False, align=PP_ALIGN.LEFT, space=6):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return p


def accent_bar(slide):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()


# ---------- Slide 1: Title ----------
s = prs.slides.add_slide(BLANK); bg(s, NAVY)
band = s.shapes.add_shape(1, 0, Inches(2.35), SW, Inches(2.9))
band.fill.solid(); band.fill.fore_color.rgb = DEEP_BLUE; band.line.fill.background()
tf = box(s, 0.9, 2.55, 11.5, 2.5)
para(tf, 'What’s an "AI Model"?', 46, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)
para(tf, "Explained Simply — in about 3 minutes", 26, ICE, align=PP_ALIGN.CENTER)
tf2 = box(s, 0.9, 5.55, 11.5, 1.2)
para(tf2, "AI Explained Simply · Episode 2  •  Anchor example: the Model dropdown in Mistral’s console",
     16, TEAL, first=True, align=PP_ALIGN.CENTER)
para(tf2, "One core AI term · plain words · everyday analogies", 14, MUTED, align=PP_ALIGN.CENTER)

# ---------- Slide 2: The un-magic secret ----------
s = prs.slides.add_slide(BLANK); bg(s, ICE); accent_bar(s)
tf = box(s, 0.7, 0.5, 12, 1.0)
para(tf, "The un-magic secret (What & Why)", 34, NAVY, bold=True, first=True)
tf = box(s, 0.9, 1.7, 11.6, 5.2)
para(tf, "A model = remembered DATA + computing MUSCLE + LANGUAGE skill", 24, DEEP_BLUE, bold=True, first=True, space=14)
for t in [
    "▪  Data = your hard drive of files — but squeezed into patterns, not kept as files",
    "▪  Muscle = the same GPUs (graphics cards) that run your video games",
    "▪  Language = a fluent, multilingual talker bolted on top",
    "▪  Brilliant at MANAGING & EXTRACTING what it has already seen",
]:
    para(tf, t, 20, INK, space=10)
para(tf, "→ The dreaming — the genuinely new idea — is still yours. You bring the head.",
     20, TEAL, bold=True, space=6)

# ---------- Slide 3: How it works ----------
s = prs.slides.add_slide(BLANK); bg(s, ICE); accent_bar(s)
tf = box(s, 0.7, 0.5, 12, 1.0)
para(tf, "How does it work?", 34, NAVY, bold=True, first=True)
tf = box(s, 0.9, 1.7, 11.6, 5.2)
para(tf, "It does ONE thing on repeat: guess the next word.", 24, DEEP_BLUE, bold=True, first=True, space=14)
for t in [
    '▪  "The cat sat on the ___"  →  mat (very likely), moon (not)',
    "▪  Pick one word, then guess the next — string them into a paragraph",
    "▪  The “layers of math” = millions of trained dials voting on the likely word",
    "▪  It’s your phone’s autocomplete, grown all the way up",
]:
    para(tf, t, 20, INK, space=10)
para(tf, "⚠  Confident ≠ correct. It predicts — so always fact-check what it tells you.",
     20, RGBColor(0xB0, 0x3A, 0x2E), bold=True, space=6)

# ---------- Slide 4: Where in Mistral ----------
s = prs.slides.add_slide(BLANK); bg(s, NAVY); accent_bar(s)
tf = box(s, 0.7, 0.5, 12, 1.0)
para(tf, "Where in the Mistral console?", 34, WHITE, bold=True, first=True)
tf = box(s, 0.9, 1.7, 11.6, 5.2)
para(tf, "Open console.mistral.ai → Le Chat or the API playground", 22, ICE, bold=True, first=True, space=14)
for t in [
    "▪  At the top of the screen: a dropdown menu labeled “Model”",
    "▪  Options include Mistral Large (smart), Small (fast & cheap), Medium",
    "▪  That dropdown IS this lesson made clickable — each option is a brain to borrow",
    "▪  Picking a model = choosing your print-quality setting before you print",
]:
    para(tf, t, 20, WHITE, space=12)

# ---------- Slide 5: Try it yourself ----------
s = prs.slides.add_slide(BLANK); bg(s, ICE); accent_bar(s)
tf = box(s, 0.7, 0.5, 12, 1.0)
para(tf, "Try it yourself!", 34, NAVY, bold=True, first=True)
tf = box(s, 0.9, 1.7, 11.6, 5.2)
for t in [
    "Step 1  —  Click the “Model” dropdown and pick Mistral Large or Small",
    "Step 2  —  Type: “Explain photosynthesis to a 7-year-old”",
    "Step 3  —  Watch it type word-by-word — that’s next-word guessing, live",
    "Bonus   —  Switch models, send the SAME question — different brain, different answer",
]:
    para(tf, t, 21, INK, space=12)
para(tf, "Takeaway: match the model to the job — don’t pay supercomputer prices for a to-do list.",
     21, TEAL, bold=True, space=8)
para(tf, "Next episode →  Prompt: how you actually talk to the model.", 18, MUTED, space=4)

out = "02_Model.pptx"
prs.save(out)
print("saved", out, "with", len(prs.slides._sldIdLst), "slides")
