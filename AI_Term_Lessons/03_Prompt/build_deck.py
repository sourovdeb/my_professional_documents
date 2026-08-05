#!/usr/bin/env python3
"""Build 03_Prompt.pptx for the "AI Explained Simply" series (Episode 3: Prompt).
Self-contained: only depends on python-pptx. Run: python3 build_deck.py
Palette matches Episodes 1 & 2 (Ocean Gradient).
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---- palette (Ocean Gradient, matches Episodes 1 & 2) ----
NAVY = RGBColor(0x21, 0x29, 0x5C)
DEEP_BLUE = RGBColor(0x06, 0x5A, 0x82)
TEAL = RGBColor(0x1C, 0x72, 0x93)
ICE = RGBColor(0xEA, 0xF2, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x16, 0x20, 0x2A)
MUTED = RGBColor(0x5B, 0x6B, 0x76)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame


def para(tf, text, size, color, bold=False, first=False, align=PP_ALIGN.LEFT, space=6, bullet=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space)
    r = p.add_run()
    r.text = ("•  " + text) if bullet else text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return p


def accent_bar(slide):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()


def title(slide, text):
    accent_bar(slide)
    tf = box(slide, 0.7, 0.45, 12.0, 1.0)
    para(tf, text, 34, NAVY, bold=True, first=True)


# ---------- Slide 1: Title ----------
s = prs.slides.add_slide(BLANK); bg(s, NAVY)
band = s.shapes.add_shape(1, 0, Inches(2.35), SW, Inches(2.9))
band.fill.solid(); band.fill.fore_color.rgb = DEEP_BLUE; band.line.fill.background()
tf = box(s, 0.9, 2.55, 11.5, 2.5)
para(tf, 'What’s a "Prompt"?', 46, WHITE, bold=True, first=True, align=PP_ALIGN.CENTER)
para(tf, "Explained Simply — in about 3 minutes", 26, ICE, align=PP_ALIGN.CENTER)
tf2 = box(s, 0.9, 5.55, 11.5, 1.2)
para(tf2, "AI Explained Simply · Episode 3 · anchored in the Mistral console", 18, ICE, first=True, align=PP_ALIGN.CENTER)

# ---------- Slide 2: What & Why ----------
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
title(s, "It’s a request, not a spell")
tf = box(s, 0.8, 1.7, 11.7, 5.2)
para(tf, "A prompt = what you ask + how clearly you ask it", 24, DEEP_BLUE, bold=True, first=True, space=12)
for t in [
    "The model reads your words — not your mind.",
    "No magic words. The whole skill is clarity + a little context.",
    "Vague in → average out.  Specific in → useful out.",
    "It’s a search box, grown up: type the task, it does the task.",
]:
    para(tf, t, 20, INK, bullet=True, space=10)

# ---------- Slide 3: How it works ----------
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
title(s, "How it works")
tf = box(s, 0.8, 1.7, 11.7, 5.2)
para(tf, "Your prompt becomes the model’s starting context", 24, DEEP_BLUE, bold=True, first=True, space=12)
for t in [
    "A model predicts the next word — the prompt is what it predicts from.",
    "Longer, clearer prompt = more for it to lean on.",
    "Like the opening line of a story: it steers everything after.",
    "Add one example, and every guess aims at “good” instead of average.",
]:
    para(tf, t, 20, INK, bullet=True, space=10)

# ---------- Slide 4: Where in Mistral ----------
s = prs.slides.add_slide(BLANK); bg(s, ICE)
title(s, "Where in Mistral (what you’d see)")
tf = box(s, 0.8, 1.7, 11.7, 5.2)
para(tf, "console.mistral.ai → Le Chat / API playground", 24, DEEP_BLUE, bold=True, first=True, space=12)
for t in [
    "The big message box at the bottom = the prompt.",
    "Everything you type there is the prompt.",
    "Nearby: a “System” field — a standing prompt (next episode).",
    "Type there, hit send, read the answer.",
]:
    para(tf, t, 20, INK, bullet=True, space=10)

# ---------- Slide 5: Try it yourself ----------
s = prs.slides.add_slide(BLANK); bg(s, WHITE)
title(s, "Try it yourself")
tf = box(s, 0.8, 1.7, 11.7, 5.2)
para(tf, "Step 1  —  type:  write an email", 20, INK, bold=True, first=True, space=8)
para(tf, "→ generic, probably wrong for your situation", 18, MUTED, space=14)
para(tf, "Step 2  —  type:  Write a short, friendly email to my landlord", 20, INK, bold=True, space=2)
para(tf, "asking to fix a leaking tap. Polite, 3 sentences, no slang.", 20, INK, bold=True, space=14)
para(tf, "Step 3  —  compare:  same model, only the prompt changed", 20, INK, bold=True, space=14)
para(tf, "Takeaway: clarity is the whole skill — it makes every AI tool better.", 22, TEAL, bold=True, space=6)

prs.save("03_Prompt.pptx")
print("Saved 03_Prompt.pptx")
