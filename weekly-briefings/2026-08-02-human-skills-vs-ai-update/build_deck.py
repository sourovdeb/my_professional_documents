#!/usr/bin/env python3
"""Build the 'Pincer Closes on Juniors' doodle deck — Aug 2026 update."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# palette (matches the July edition)
INK   = RGBColor(0x22, 0x22, 0x2B)
PAPER = RGBColor(0xFB, 0xF7, 0xEF)
ACCENT= RGBColor(0xE8, 0x55, 0x3A)  # warm orange-red
BLUE  = RGBColor(0x2F, 0x6F, 0xB0)
GREEN = RGBColor(0x3B, 0x8C, 0x5A)
GREY  = RGBColor(0x8A, 0x86, 0x7E)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height

def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s

def txt(s, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font="Comic Sans MS", line_spacing=1.05):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = font
    return tb

def doodle(s, shape, x, y, w, h, color=INK, fill=None, wdt=2.5):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = color; sp.line.width = Pt(wdt)
    sp.shadow.inherit = False
    return sp

def hbar(s, x, y, w_in, label, val, color, maxw=6.2):
    doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(w_in), Inches(0.52),
           color=INK, fill=color, wdt=2)
    txt(s, Emu(int(x + Inches(w_in) + Inches(0.1))), y-Inches(0.02), Inches(2.2),
        Inches(0.55), val, 16, INK, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    txt(s, x-Inches(3.4), y-Inches(0.02), Inches(3.25), Inches(0.55), label, 13,
        GREY, False, PP_ALIGN.RIGHT, MSO_ANCHOR.MIDDLE)

# ---------- Slide 1: title ----------
s = slide()
txt(s, Inches(0.9), Inches(1.25), Inches(11.7), Inches(1.6),
    "The Pincer Closes on Juniors", 54, INK, True, PP_ALIGN.LEFT)
doodle(s, MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(2.6), Inches(7.4), Inches(0.16),
       color=ACCENT, fill=ACCENT, wdt=1)
txt(s, Inches(0.95), Inches(2.9), Inches(11.6), Inches(1.0),
    "Human skills vs AI — the August 2026 update.", 28, ACCENT, True)
txt(s, Inches(0.95), Inches(3.75), Inches(11.6), Inches(1.0),
    "Sentiment → mechanism.  We now see WHERE it lands (juniors)\nand HOW FAST people adapt (+185% critical-thinking enrollments).",
    20, GREY)
txt(s, Inches(0.95), Inches(5.3), Inches(11.6), Inches(0.7),
    "5-minute read  ·  deltas vs the 19 July baseline  ·  2 August 2026", 18, GREY)
doodle(s, MSO_SHAPE.LIGHTNING_BOLT, Inches(10.9), Inches(0.9), Inches(1.4), Inches(2.0),
       color=ACCENT, fill=RGBColor(0xFF,0xE7,0x9c), wdt=3)

# ---------- Slide 2: the one line ----------
s = slide()
txt(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.8), "The one line", 34, ACCENT, True)
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.6), Inches(11.3), Inches(3.9),
       color=INK, fill=RGBColor(0xFF,0xF3,0xD9), wdt=3)
txt(s, Inches(1.5), Inches(1.95), Inches(10.3), Inches(3.3),
    "Two weeks ago this monitor reported FEAR.\nThis edition reports DATA —\nand it lands hardest on early-career workers.\n\n"
    "The AI-skills premium held. But the entry-level door\ndidn't slam — it NARROWED and RAISED its bar:\n"
    "surviving junior jobs now demand senior judgment on day one.", 24, INK, True,
    PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
txt(s, Inches(1.0), Inches(5.85), Inches(11.3), Inches(1.0),
    "The July thesis held:  AI × human, not AI vs human.  What's new is the COORDINATE.",
    22, ACCENT, True, PP_ALIGN.CENTER)

# ---------- Slide 3: the shift (before/after) ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "The shift: from fear to mechanism", 32, ACCENT, True)
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.6), Inches(5.4), Inches(4.9),
       color=GREY, fill=RGBColor(0xEF,0xEC,0xE4), wdt=2.5)
txt(s, Inches(1.2), Inches(1.8), Inches(4.8), Inches(0.6), "19 JULY — fear", 22, GREY, True)
txt(s, Inches(1.2), Inches(2.6), Inches(4.9), Inches(3.7),
    "“60% expect AI to\ndestroy more jobs.”\n\n“22% feel their\njob is safe.”\n\n“Learn durable\nskills.”  (advice)",
    20, INK)
doodle(s, MSO_SHAPE.RIGHT_ARROW, Inches(6.45), Inches(3.6), Inches(0.9), Inches(0.9),
       color=ACCENT, fill=ACCENT, wdt=1)
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.6), Inches(5.0), Inches(4.9),
       color=BLUE, fill=RGBColor(0xE6,0xEF,0xF7), wdt=3)
txt(s, Inches(7.8), Inches(1.8), Inches(4.5), Inches(0.6), "2 AUGUST — mechanism", 22, BLUE, True)
txt(s, Inches(7.8), Inches(2.6), Inches(4.5), Inches(3.7),
    "Stanford AI Index:\n−20% employment for\ndevs aged 22–25.\n\nEntry roles “seniorised”:\n+35% ask for judgment.\n\n"
    "Critical-thinking\nenrollments +185%.\nPeople are ACTING.", 19, INK)

# ---------- Slide 4: premium held (numbers) ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "The premium held — the primary source got sharper", 30, ACCENT, True)
txt(s, Inches(0.9), Inches(1.4), Inches(11.5), Inches(0.6),
    "PwC 2026 Barometer (1bn+ ads, 27 countries) — now direct, not secondhand:", 17, GREY, True)
hbar(s, Inches(4.4), Inches(2.35), 2.6, "AI wage premium", "62%  (held)", ACCENT)
hbar(s, Inches(4.4), Inches(3.15), 4.8, "AI-job demand growth", "+69%  vs +9% mkt (8×)", BLUE)
hbar(s, Inches(4.4), Inches(3.95), 5.6, "AI-skill demand YoY", "+144%", BLUE)
hbar(s, Inches(4.4), Inches(4.75), 5.9, "Productivity, top firms", "163% gain", GREEN)
hbar(s, Inches(4.4), Inches(5.55), 1.2, "Premium — govt floor", "16%", GREY)
hbar(s, Inches(4.4), Inches(6.10), 6.0, "Premium — consumer top", "118%", GREY)
txt(s, Inches(0.9), Inches(6.75), Inches(11.5), Inches(0.6),
    "Read: not a bubble popping — the money held while the VOLUME accelerated to 8× the wider market.",
    16, INK, True, PP_ALIGN.CENTER)

# ---------- Slide 5: the entry-level squeeze ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "Durable skills — now with a body count at the bottom rung", 28, ACCENT, True)
facts = [
    ("−20%", "software-dev employment,\nages 22–25, since 2024\n(Stanford AI Index)", ACCENT),
    ("7×", "more likely an entry role\nnow demands SENIOR\njudgment (PwC)", BLUE),
    ("21%", "of firms already FROZE\nentry-level hiring; ~half\nexpect to end it by 2027", GREEN),
]
x = Inches(0.9)
for val, body, col in facts:
    doodle(s, MSO_SHAPE.OVAL, x, Inches(1.75), Inches(3.7), Inches(2.7),
           color=col, fill=PAPER, wdt=3)
    txt(s, x, Inches(2.25), Inches(3.7), Inches(0.9), val, 40, col, True, PP_ALIGN.CENTER)
    txt(s, x, Inches(3.2), Inches(3.7), Inches(1.1), body, 14, INK, False, PP_ALIGN.CENTER)
    x = Emu(int(x + Inches(4.1)))
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(4.9), Inches(11.5), Inches(1.9),
       color=INK, fill=RGBColor(0xFF,0xF3,0xD9), wdt=2.5)
txt(s, Inches(1.2), Inches(5.05), Inches(11.0), Inches(1.7),
    "THE NUANCE THAT MATTERS: this is a HIRING freeze, not a firing wave.\n"
    "AI is suppressing the FRONT DOOR, not eliminating incumbents. The danger isn't\n"
    "“everyone gets replaced” — it's “no one gets the first rung.”", 18, INK, True,
    PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

# ---------- Slide 6: the sawn-off ladder example ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "What it looks like: the sawn-off ladder", 30, ACCENT, True)
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.6), Inches(5.4), Inches(4.9),
       color=GREY, fill=RGBColor(0xEF,0xEC,0xE4), wdt=2.5)
txt(s, Inches(1.2), Inches(1.8), Inches(4.8), Inches(0.6), "2024 junior", 22, GREY, True)
txt(s, Inches(1.2), Inches(2.6), Inches(4.9), Inches(3.7),
    "Year one = cleaning\ndata, reconciling\nspreadsheets.\n\nBoring — but it's how\nyou learn what\n“wrong” looks like.\n\n"
    "You climbed rung by\nrung.", 19, INK)
doodle(s, MSO_SHAPE.RIGHT_ARROW, Inches(6.45), Inches(3.6), Inches(0.9), Inches(0.9),
       color=ACCENT, fill=ACCENT, wdt=1)
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.6), Inches(5.0), Inches(4.9),
       color=BLUE, fill=RGBColor(0xE6,0xEF,0xF7), wdt=3)
txt(s, Inches(7.8), Inches(1.8), Inches(4.5), Inches(0.6), "2026 junior", 22, BLUE, True)
txt(s, Inches(7.8), Inches(2.6), Inches(4.5), Inches(3.7),
    "An agent does rungs\n1 & 2 in seconds.\n\nThe job ad now reads:\n“exercise judgment\nover AI outputs,\npresent to\nstakeholders.”\n\n"
    "You're asked to START\nat rung three.", 19, INK)
txt(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.7),
    "Winners build rung-three skills BEFORE they're hired.", 20, ACCENT, True, PP_ALIGN.CENTER)

# ---------- Slide 7: where the learning energy goes ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "People are voting with their course tabs", 30, GREEN, True)
txt(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(0.6),
    "Coursera Job Skills Report 2026 — YoY enrollment growth:", 18, GREY, True)
bars = [
    ("GenAI (record: 14/min)", 234, ACCENT),
    ("Critical thinking (GenAI)", 185, BLUE),
    ("Critical thinking (Data)", 168, BLUE),
    ("Critical thinking (IT)", 101, GREEN),
]
y = Inches(2.2)
maxpct = 234
for label, pct, col in bars:
    w_in = 0.9 + (pct / maxpct) * 8.5
    doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.7), y, Inches(w_in), Inches(0.62),
           color=INK, fill=col, wdt=2)
    txt(s, Emu(int(Inches(3.7) + Inches(w_in) + Inches(0.1))), y, Inches(1.6),
        Inches(0.62), f"+{pct}%", 18, INK, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    txt(s, Inches(0.4), y, Inches(3.15), Inches(0.62), label, 14, GREY, False,
        PP_ALIGN.RIGHT, MSO_ANCHOR.MIDDLE)
    y = Emu(int(y + Inches(0.82)))
txt(s, Inches(0.9), Inches(5.75), Inches(11.5), Inches(1.4),
    "Coursera's own framing matches this monitor exactly:\n"
    "the human role is moving “from collaborator to EXPERT VALIDATOR of the output.”\n"
    "Skilliday deepened too: Hilton — 72% would take work time off for a skill-trip.",
    17, INK, True, PP_ALIGN.CENTER)

# ---------- Slide 8: the tool ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "🛠 A tool to try this week", 32, ACCENT, True)
txt(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(0.6),
    "The “Rung-Three Résumé” test  (10 minutes, early-career)", 22, INK, True)
steps = [
    "1.  Write the job you want. List 3 tasks in it an AI agent now does end-to-end. Cross them out.",
    "2.  For each, name the judgment on top: is it right? is it wise? should we act? — that's rung three.",
    "3.  Score yourself 0–2 on each judgment call. Zeros are your gap list.",
    "4.  Manufacture reps: run real data through an AI, write “here's what it got wrong.” Do 5 = portfolio.",
    "5.  Habit: ask any AI to argue the OPPOSITE of your conclusion. Autopilot → critical-thinking gym.",
]
y = Inches(2.15)
for st in steps:
    doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), y, Inches(11.5), Inches(0.74),
           color=BLUE, fill=RGBColor(0xF0,0xF5,0xFB), wdt=2)
    txt(s, Inches(1.15), y+Inches(0.06), Inches(11.0), Inches(0.64), st, 16, INK, False,
        PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    y = Emu(int(y + Inches(0.87)))
txt(s, Inches(0.9), Inches(6.65), Inches(11.5), Inches(0.7),
    "The market stopped buying “I can produce the output.” It buys “I can be trusted to VALIDATE it.”",
    16, ACCENT, True, PP_ALIGN.CENTER)

# ---------- Slide 9: scoreboard / sources ----------
s = slide()
txt(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.9),
    "The scoreboard vs 19 July", 34, ACCENT, True)
rows = [
    ("PwC AI wage premium", "62% → 62%", "held", GREY),
    ("AI-job demand growth", "→ +69% (8×)", "new", BLUE),
    ("Entry-level dev employment (22–25)", "fear → −20%", "HARD DATA", ACCENT),
    ("Critical-thinking enrollments", "→ +185% YoY", "new", GREEN),
    ("Skilliday intent (Gen Z)", "57% plan → 72% off-work", "deepened", GREEN),
]
y = Inches(1.55)
for label, move, tag, col in rows:
    doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), y, Inches(11.5), Inches(0.66),
           color=col, fill=PAPER, wdt=2)
    txt(s, Inches(1.1), y, Inches(6.0), Inches(0.66), label, 16, INK, False,
        PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    txt(s, Inches(7.0), y, Inches(3.4), Inches(0.66), move, 16, INK, True,
        PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    txt(s, Inches(10.4), y, Inches(1.9), Inches(0.66), tag, 15, col, True,
        PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    y = Emu(int(y + Inches(0.78)))
txt(s, Inches(0.9), Inches(5.65), Inches(11.5), Inches(1.6),
    "Biggest shift: the entry-level story went from worker FEAR to measured CONTRACTION.\n\n"
    "Sources: PwC 2026 AI Jobs Barometer · SUCCESS · Coursera Job Skills Report 2026 · No Jitter ·\n"
    "Stanford HAI AI Index · Goldman Sachs · Axis Intelligence · Hilton/Priority Pass skillcation.\n"
    "(Full links in the accompanying README.)", 15, GREY, False, PP_ALIGN.CENTER)

import os
here = os.path.dirname(os.path.abspath(__file__))
out = os.path.join(here, "Human_Skills_vs_AI_Update_Aug2026.pptx")
prs.save(out)
print("saved deck with", len(prs.slides._sldIdLst), "slides ->", out)
