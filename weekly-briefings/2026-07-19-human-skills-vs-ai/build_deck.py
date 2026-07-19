#!/usr/bin/env python3
"""Build the 'Human-Skills Premium 2026' doodle deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# palette
INK   = RGBColor(0x22, 0x22, 0x2B)
PAPER = RGBColor(0xFB, 0xF7, 0xEF)
ACCENT= RGBColor(0xE8, 0x55, 0x3A)  # warm orange-red
BLUE  = RGBColor(0x2F, 0x6F, 0xB0)
GREEN = RGBColor(0x3B, 0x8C, 0x5A)
GREY  = RGBColor(0x8A, 0x86, 0x7E)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height

def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s

def txt(s, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font="Comic Sans MS", line_spacing=1.05):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = font
    return tb

def doodle(s, shape, x, y, w, h, color=INK, fill=None, wdt=2.5):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = color; sp.line.width = Pt(wdt)
    sp.shadow.inherit = False
    return sp

def blob_label(s, x, y, w, h, label, num, ncolor):
    c = doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, color=ncolor, fill=PAPER, wdt=3)
    doodle(s, MSO_SHAPE.OVAL, x-Inches(0.28), y-Inches(0.28), Inches(0.85), Inches(0.85),
           color=INK, fill=ncolor, wdt=2.5)
    txt(s, x-Inches(0.28), y-Inches(0.30), Inches(0.85), Inches(0.85), num, 26, PAPER,
        True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    return c

# ---------- Slide 1: title ----------
s = slide()
txt(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.6),
    "The Human-Skills Premium", 60, INK, True, PP_ALIGN.LEFT)
doodle(s, MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(2.95), Inches(6.6), Inches(0.16),
       color=ACCENT, fill=ACCENT, wdt=1)
txt(s, Inches(0.95), Inches(3.25), Inches(11.4), Inches(1.0),
    "What to learn when AI learns everything.", 30, ACCENT, True)
txt(s, Inches(0.95), Inches(4.15), Inches(11.4), Inches(1.0),
    "2026 mid-year trend brief  ·  5-minute read  ·  19 July 2026", 20, GREY)
# doodle: lightbulb + gear
doodle(s, MSO_SHAPE.LIGHTNING_BOLT, Inches(10.7), Inches(1.2), Inches(1.4), Inches(2.0),
       color=ACCENT, fill=RGBColor(0xFF,0xE7,0x9c), wdt=3)
doodle(s, MSO_SHAPE.SUN, Inches(9.5), Inches(4.6), Inches(1.7), Inches(1.7),
       color=BLUE, fill=None, wdt=3)
txt(s, Inches(9.5), Inches(5.15), Inches(1.7), Inches(0.7), "AI ×\nHUMAN", 16, BLUE,
    True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

# ---------- Slide 2: the one line ----------
s = slide()
txt(s, Inches(0.9), Inches(0.6), Inches(11.5), Inches(0.8), "The one line", 34, ACCENT, True)
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.7), Inches(11.3), Inches(3.4),
       color=INK, fill=RGBColor(0xFF,0xF3,0xD9), wdt=3)
txt(s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(2.9),
    "AI is not erasing the value of human skills —\nit is REPRICING them.\n\n"
    "The tasks it automates get cheaper.\nThe judgment, taste & trust it can't fake get "
    "more expensive.", 30, INK, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
txt(s, Inches(1.0), Inches(5.5), Inches(11.3), Inches(1.2),
    "The 2026 move isn't “AI vs. human.”  It's  AI × human.", 26, ACCENT, True, PP_ALIGN.CENTER)

# ---------- Slide 3: three forces ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "Three forces on your career right now", 32, ACCENT, True)
cols = [
    (Inches(0.9),  "1", BLUE,  "DURABLE HUMAN SKILLS",
     "Critical thinking is the\n#1 skill leaders want\n(73%).\n\nJudgment · leadership ·\nEQ · human-in-the-loop"),
    (Inches(5.05), "2", ACCENT,"AI JOB-MARKET SHIFTS",
     "+62% wage premium\n(PwC, up from 57%).\n\nBut skills now expire\nin 2.5–5 years."),
    (Inches(9.2),  "3", GREEN, "NICHE / LIVED LEARNING",
     "The “skilliday” boom.\n\nLanguages, crafts,\nwellness — half of\nEurope, led by Gen Z."),
]
for x, n, col, head, body in cols:
    blob_label(s, x, Inches(1.9), Inches(3.55), Inches(3.9), head, n, col)
    txt(s, x+Inches(0.2), Inches(2.25), Inches(3.15), Inches(0.9), head, 17, col, True,
        PP_ALIGN.CENTER)
    txt(s, x+Inches(0.25), Inches(3.15), Inches(3.05), Inches(2.5), body, 15, INK)
# arrow to meta-skill
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.0), Inches(6.2), Inches(7.3), Inches(0.9),
       color=INK, fill=RGBColor(0xE7,0xF0,0xE7), wdt=2.5)
txt(s, Inches(3.0), Inches(6.25), Inches(7.3), Inches(0.8),
    "META-SKILL underneath all three: learning how to learn", 20, GREEN, True,
    PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

# ---------- Slide 4: the numbers (doodle bar chart) ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "Two price tags: the carrot & the stick", 32, ACCENT, True)
txt(s, Inches(0.9), Inches(1.5), Inches(5.6), Inches(0.6), "THE CARROT — AI wage premium", 18, BLUE, True)
# bars
base_y = Inches(4.6)
def bar(x, h_in, label, val, color):
    doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Emu(int(base_y - Inches(h_in))),
           Inches(1.1), Inches(h_in), color=INK, fill=color, wdt=2)
    txt(s, x-Inches(0.15), Emu(int(base_y - Inches(h_in) - Inches(0.5))), Inches(1.4),
        Inches(0.5), val, 18, INK, True, PP_ALIGN.CENTER)
    txt(s, x-Inches(0.15), base_y, Inches(1.4), Inches(0.7), label, 12, GREY, False, PP_ALIGN.CENTER)
bar(Inches(1.2), 1.2, "Lightcast", "+28%", BLUE)
bar(Inches(3.0), 2.6, "PwC 2026", "+62%", ACCENT)
txt(s, Inches(0.9), Inches(5.4), Inches(5.6), Inches(1.6),
    "…and it's WIDENING: 57% → 62% in one year.\n51% of AI-skill jobs are now OUTSIDE IT.",
    16, INK)
# stick side
doodle(s, MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.5), Inches(0.03), Inches(5.2),
       color=GREY, fill=GREY, wdt=1)
txt(s, Inches(7.2), Inches(1.5), Inches(5.6), Inches(0.6), "THE STICK — skills expire", 18, ACCENT, True)
doodle(s, MSO_SHAPE.PIE, Inches(7.4), Inches(2.3), Inches(2.0), Inches(2.0),
       color=INK, fill=RGBColor(0xFF,0xD9,0xD0), wdt=2.5)
txt(s, Inches(7.4), Inches(4.35), Inches(2.0), Inches(0.5), "half-life\n2.5–5 yrs", 14, INK, True, PP_ALIGN.CENTER)
txt(s, Inches(9.8), Inches(2.4), Inches(3.2), Inches(3.0),
    "• 70% of niche tech\n  knowledge stale in\n  18 months.\n\n• 39% of core skills\n  obsolete by 2030\n  (WEF).", 16, INK)

# ---------- Slide 5: the human-in-the-loop example ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "What the work actually looks like now", 32, ACCENT, True)
# before box
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.7), Inches(5.4), Inches(4.4),
       color=GREY, fill=RGBColor(0xEF,0xEC,0xE4), wdt=2.5)
txt(s, Inches(1.2), Inches(1.9), Inches(4.8), Inches(0.6), "BEFORE", 22, GREY, True)
txt(s, Inches(1.2), Inches(2.6), Inches(4.8), Inches(3.3),
    "Analyst spends Monday\nbuilding the dashboard\nby hand.\n\nValue = producing\nthe output.", 20, INK)
# arrow
doodle(s, MSO_SHAPE.RIGHT_ARROW, Inches(6.45), Inches(3.4), Inches(0.9), Inches(0.9),
       color=ACCENT, fill=ACCENT, wdt=1)
# after box
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.7), Inches(5.0), Inches(4.4),
       color=BLUE, fill=RGBColor(0xE6,0xEF,0xF7), wdt=3)
txt(s, Inches(7.8), Inches(1.9), Inches(4.4), Inches(0.6), "NOW", 22, BLUE, True)
txt(s, Inches(7.8), Inches(2.6), Inches(4.4), Inches(3.3),
    "AI drafts it in 4 min.\n\nShe spots the tracking\nbug, reframes it for the\nexec, decides what NOT\n"
    "to recommend.\n\nValue = the judgment.", 20, INK)
txt(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.8),
    "Same salary line. Completely different skill doing the earning.", 20, ACCENT, True, PP_ALIGN.CENTER)

# ---------- Slide 6: skillidays ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "People are voting with their holidays: the “skilliday”", 30, GREEN, True)
doodle(s, MSO_SHAPE.SUN, Inches(10.6), Inches(0.4), Inches(1.9), Inches(1.9),
       color=RGBColor(0xE0,0xA5,0x2A), fill=RGBColor(0xFF,0xE7,0x9c), wdt=3)
stats = [
    ("~50%", "of Europeans plan to\nlearn a skill on holiday\n(Mastercard, 27k people)"),
    ("57%", "of 18–24s are planning\na skill-based trip"),
    ("1 in 3", "of Gen Z prefer a skill\nover a souvenir"),
    ("76%", "say a skill outlasts\nany souvenir"),
]
x = Inches(0.9)
for val, body in stats:
    doodle(s, MSO_SHAPE.OVAL, x, Inches(2.1), Inches(2.7), Inches(2.7),
           color=GREEN, fill=PAPER, wdt=3)
    txt(s, x, Inches(2.7), Inches(2.7), Inches(0.9), val, 34, GREEN, True, PP_ALIGN.CENTER)
    txt(s, x, Inches(3.55), Inches(2.7), Inches(1.1), body, 13, INK, False, PP_ALIGN.CENTER)
    x = Emu(int(x + Inches(3.0)))
txt(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.6),
    "Top categories: languages · wellness & movement (25%) · crafts (24%) · creative arts (23%).\n\n"
    "Why it matters for work: pottery and negotiation are cousins — tactile, feedback-driven,\n"
    "and impossible to fully outsource to a model.", 17, INK)

# ---------- Slide 7: the formula ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "The actionable combination", 32, ACCENT, True)
txt(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(0.6),
    "Isolated topics are out. Combinations are in.", 20, GREY, True)
boxes = [("AI\nleverage", BLUE), ("CRITICAL\nTHINKING", ACCENT), ("DOMAIN\nEXPERTISE", GREEN)]
x = Inches(1.1)
for i, (label, col) in enumerate(boxes):
    doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.6), Inches(3.0), Inches(2.2),
           color=col, fill=PAPER, wdt=4)
    txt(s, x, Inches(3.1), Inches(3.0), Inches(1.2), label, 24, col, True,
        PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    if i < 2:
        txt(s, Emu(int(x+Inches(3.0))), Inches(2.9), Inches(0.9), Inches(1.6), "×", 54, INK,
            True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    x = Emu(int(x + Inches(3.9)))
txt(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(1.4),
    "= a 2026-proof role.\nMaster the meta-skill (learning how to learn) and you don't have to\n"
    "predict which tool wins — you just re-tool faster than the person next to you.",
    20, INK, True, PP_ALIGN.CENTER)

# ---------- Slide 8: the tool ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "🛠 A tool to try this week", 32, ACCENT, True)
txt(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(0.6),
    "The 15-Minute Skill-Portfolio Canvas", 24, INK, True)
steps = [
    "1.  List the 5 tasks that fill most of your week.",
    "2.  Tag each:  🤖 AI can do it  ·  🧠 needs my judgment  ·  🎨 needs a human touch.",
    "3.  Count your 🤖 tags — that's your automation exposure.",
    "4.  For each 🤖, write the downstream judgment it frees you for (validate, decide, persuade).",
    "5.  Grow ONE durable + ONE refreshable skill this quarter. Review in 90 days — beat the half-life.",
]
y = Inches(2.2)
for st in steps:
    doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), y, Inches(11.5), Inches(0.72),
           color=BLUE, fill=RGBColor(0xF0,0xF5,0xFB), wdt=2)
    txt(s, Inches(1.15), y+Inches(0.06), Inches(11.0), Inches(0.62), st, 17, INK, False,
        PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    y = Emu(int(y + Inches(0.85)))
txt(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.8),
    "Bonus habit: ask an AI to argue the OPPOSITE of your conclusion. Autopilot → critical-thinking gym.",
    16, ACCENT, True, PP_ALIGN.CENTER)

# ---------- Slide 9: closing / sources ----------
s = slide()
txt(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(1.0),
    "The takeaway", 40, ACCENT, True)
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.9), Inches(11.3), Inches(2.3),
       color=INK, fill=RGBColor(0xFF,0xF3,0xD9), wdt=3)
txt(s, Inches(1.4), Inches(2.1), Inches(10.5), Inches(1.9),
    "People are learning to stay relevant in an AI-augmented world —\n"
    "and the meta-skill is learning how to learn, fast.", 26, INK, True,
    PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.9), Inches(4.5), Inches(11.5), Inches(2.2),
    "Sources: Korn Ferry TA Trends 2026 · Lightcast · PwC 2026 Global AI Jobs Barometer ·\n"
    "ADP Research · Mastercard/Euronews skilliday survey · IBM skill half-life · WEF Future of Jobs ·\n"
    "Stanford HAI AI Index 2026 · Toggl meta-skills.  (Full links in the accompanying README.)",
    15, GREY, False, PP_ALIGN.CENTER)

prs.save("/home/user/my_professional_documents/weekly-briefings/2026-07-19-human-skills-vs-ai/Human_Skills_Premium_2026.pptx")
print("saved deck with", len(prs.slides._sldIdLst), "slides")
