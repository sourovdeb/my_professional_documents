#!/usr/bin/env python3
"""Build the 'Two-Track Labour Market 2026' doodle deck (trend brief #2, 24 Jul)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# palette
INK   = RGBColor(0x22, 0x22, 0x2B)
PAPER = RGBColor(0xFB, 0xF7, 0xEF)
ACCENT= RGBColor(0xE8, 0x55, 0x3A)  # warm orange-red
BLUE  = RGBColor(0x2F, 0x6F, 0xB0)
GREEN = RGBColor(0x3B, 0x8C, 0x5A)
GREY  = RGBColor(0x8A, 0x86, 0x7E)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height

def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s

def txt(s, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font="Comic Sans MS", line_spacing=1.05):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = font
    return tb

def doodle(s, shape, x, y, w, h, color=INK, fill=None, wdt=2.5):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = color; sp.line.width = Pt(wdt)
    sp.shadow.inherit = False
    return sp

# ---------- Slide 1: title ----------
s = slide()
txt(s, Inches(0.9), Inches(1.35), Inches(11.9), Inches(1.6),
    "The Two-Track Labour Market", 54, INK, True, PP_ALIGN.LEFT)
doodle(s, MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(2.75), Inches(7.4), Inches(0.16),
       color=ACCENT, fill=ACCENT, wdt=1)
txt(s, Inches(0.95), Inches(3.05), Inches(11.6), Inches(1.0),
    "The AI-skills story just split in two.", 30, ACCENT, True)
txt(s, Inches(0.95), Inches(3.95), Inches(11.6), Inches(1.0),
    "2026 trend brief #2  ·  5-minute read  ·  24 July 2026  ·  follow-up to 19 Jul", 19, GREY)
# doodle: forking arrows
doodle(s, MSO_SHAPE.UP_RIBBON, Inches(10.6), Inches(1.1), Inches(1.9), Inches(1.3),
       color=BLUE, fill=RGBColor(0xE6,0xEF,0xF7), wdt=3)
txt(s, Inches(10.6), Inches(1.45), Inches(1.9), Inches(0.7), "+118%", 20, BLUE, True,
    PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
doodle(s, MSO_SHAPE.DOWN_RIBBON, Inches(10.6), Inches(4.9), Inches(1.9), Inches(1.3),
       color=ACCENT, fill=RGBColor(0xFF,0xD9,0xD0), wdt=3)
txt(s, Inches(10.6), Inches(5.25), Inches(1.9), Inches(0.7), "+16%", 20, ACCENT, True,
    PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

# ---------- Slide 2: the one line ----------
s = slide()
txt(s, Inches(0.9), Inches(0.6), Inches(11.5), Inches(0.8), "The one line", 34, ACCENT, True)
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.7), Inches(11.3), Inches(3.6),
       color=INK, fill=RGBColor(0xFF,0xF3,0xD9), wdt=3)
txt(s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(3.1),
    "Last brief's story was REPRICING.\n\n"
    "Five days of fresh data sharpen it into something colder:\nBIFURCATION.\n\n"
    "The same forces are splitting the workforce into TWO TRACKS —\n"
    "and the line is who gets to do the judgment.", 26, INK, True, PP_ALIGN.LEFT,
    MSO_ANCHOR.MIDDLE)
txt(s, Inches(1.0), Inches(5.6), Inches(11.3), Inches(1.2),
    "The rung it's quietly kicking away is entry-level.", 24, ACCENT, True, PP_ALIGN.CENTER)

# ---------- Slide 3: what changed (delta table) ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "What changed since 19 July", 32, ACCENT, True)
rows = [
    ("PwC wage premium", "62% average", "62% avg — but 118% consumer / 16% public", BLUE),
    ("AI-job demand", "“growing fast”", "+69% vs +9% overall  (~8x faster)", GREEN),
    ("Durable-skill demand", "73% of leaders ask", "76% of ALL postings want >=1 durable skill", ACCENT),
    ("Entry-level", "not tracked", "seniorised +35% / ordinary entry -10%", INK),
]
y = Inches(1.6)
for label, old, new, col in rows:
    doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), y, Inches(11.5), Inches(1.15),
           color=col, fill=PAPER, wdt=2.5)
    txt(s, Inches(1.1), y+Inches(0.12), Inches(2.9), Inches(0.9), label, 17, col, True,
        PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.0), y+Inches(0.12), Inches(2.7), Inches(0.9), "was:\n" + old, 13, GREY,
        False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    txt(s, Inches(6.8), y+Inches(0.12), Inches(5.4), Inches(0.9), "now:  " + new, 15, INK,
        True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    y = Emu(int(y + Inches(1.32)))
txt(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.5),
    "The headline numbers held. What's new is the SHAPE.", 17, ACCENT, True, PP_ALIGN.CENTER)

# ---------- Slide 4: the premium is a fan (bar chart) ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "Force 1: the premium is a fan, not a number", 30, ACCENT, True)
txt(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(0.6),
    "AI-skills wage premium by sector (PwC 2026)", 18, GREY, True)
base_y = Inches(5.9)
def bar(x, h_in, label, val, color):
    doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Emu(int(base_y - Inches(h_in))),
           Inches(1.5), Inches(h_in), color=INK, fill=color, wdt=2)
    txt(s, x-Inches(0.25), Emu(int(base_y - Inches(h_in) - Inches(0.55))), Inches(2.0),
        Inches(0.5), val, 22, INK, True, PP_ALIGN.CENTER)
    txt(s, x-Inches(0.25), base_y, Inches(2.0), Inches(0.9), label, 13, GREY, False, PP_ALIGN.CENTER)
bar(Inches(1.4), 4.0, "consumer /\ncustomer-facing", "118%", BLUE)
bar(Inches(4.6), 2.1, "the “average”\neveryone quotes", "62%", GREY)
bar(Inches(7.8), 0.55, "government /\npublic sector", "16%", ACCENT)
txt(s, Inches(9.9), Inches(2.2), Inches(3.1), Inches(3.5),
    "Where the AI\nmeets a paying\ncustomer, the\npremium explodes.\n\n"
    "It pays for the\nJUDGMENT you wrap\naround the tool —\nnot the tool.", 16, INK, True)

# ---------- Slide 5: professionalise vs democratise ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "Force 2: professionalising AND democratising at once", 28, ACCENT, True)
# left track
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.7), Inches(5.4), Inches(4.4),
       color=BLUE, fill=RGBColor(0xE6,0xEF,0xF7), wdt=3.5)
txt(s, Inches(1.2), Inches(1.9), Inches(4.8), Inches(0.7), "PROFESSIONALISED", 22, BLUE, True)
txt(s, Inches(1.2), Inches(2.55), Inches(4.8), Inches(0.6), "“force multiplier for experts”", 15, GREY, True)
txt(s, Inches(1.2), Inches(3.2), Inches(4.8), Inches(2.8),
    "• 2× the job growth\n\n• salaries rising\n  42% FASTER\n\n• needs MORE human\n  skill, not less", 19, INK)
# right track
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.7), Inches(5.4), Inches(4.4),
       color=ACCENT, fill=RGBColor(0xFF,0xE7,0xDF), wdt=2.5)
txt(s, Inches(7.3), Inches(1.9), Inches(4.8), Inches(0.7), "DEMOCRATISED", 22, ACCENT, True)
txt(s, Inches(7.3), Inches(2.55), Inches(4.8), Inches(0.6), "“anyone can now do the task”", 15, GREY, True)
txt(s, Inches(7.3), Inches(3.2), Inches(4.8), Inches(2.8),
    "• slower expansion\n\n• wage pressure\n\n• the task\n  commoditises", 19, INK)
txt(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.8),
    "Which track your role lands on is being decided NOW — by how much judgment it carries.",
    18, INK, True, PP_ALIGN.CENTER)

# ---------- Slide 6: the bottom rung ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "Force 3: the new casualty is the bottom rung", 30, ACCENT, True)
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.6), Inches(5.4), Inches(2.0),
       color=GREEN, fill=RGBColor(0xE7,0xF0,0xE7), wdt=3)
txt(s, Inches(1.1), Inches(1.75), Inches(5.0), Inches(1.8),
    "“Seniorised” entry roles\n(AI-exposed fields)\n\n+35%  since 2019", 22, GREEN, True,
    PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.6), Inches(5.4), Inches(2.0),
       color=ACCENT, fill=RGBColor(0xFF,0xD9,0xD0), wdt=3)
txt(s, Inches(7.2), Inches(1.75), Inches(5.0), Inches(1.8),
    "Ordinary entry roles\n\n−10%", 22, ACCENT, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.9), Inches(3.9), Inches(11.5), Inches(1.0),
    "Entry roles in AI-exposed fields are 7× more likely to demand SENIOR-level skills.",
    20, INK, True, PP_ALIGN.CENTER)
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.6),
       color=INK, fill=RGBColor(0xFF,0xF3,0xD9), wdt=2.5)
txt(s, Inches(1.8), Inches(5.15), Inches(9.7), Inches(1.35),
    "AI now does the easy tasks — so the market deleted the “learn on the job slowly” phase.\n"
    "Arrive already carrying a durable skill and a validated judgment call.", 18, INK, True,
    PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

# ---------- Slide 7: durable skills hard numbers ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "Force 4: durable skills, now with hard numbers", 29, GREEN, True)
txt(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(0.6),
    "Lightcast — 75M+ job postings analysed", 17, GREY, True)
stats = [
    ("76%", "of ALL postings request\n≥ 1 durable skill"),
    ("47%", "request 3+ durable skills\n(up 13 pts since 2021)"),
    ("8/10", "of the top-10 skills\nare durable, not technical"),
    ("#1", "critical thinking =\nhighest-value AI complement\n(McKinsey)"),
]
x = Inches(0.9)
for val, body in stats:
    doodle(s, MSO_SHAPE.OVAL, x, Inches(2.2), Inches(2.7), Inches(2.7),
           color=GREEN, fill=PAPER, wdt=3)
    txt(s, x, Inches(2.85), Inches(2.7), Inches(0.9), val, 34, GREEN, True, PP_ALIGN.CENTER)
    txt(s, x, Inches(3.7), Inches(2.7), Inches(1.1), body, 13, INK, False, PP_ALIGN.CENTER)
    x = Emu(int(x + Inches(3.0)))
txt(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.4),
    "Durable skills aren't the soft nice-to-have layer anymore.\n"
    "On the actual postings, they're the MAJORITY of what employers ask for by name.",
    19, INK, True, PP_ALIGN.CENTER)

# ---------- Slide 8: sentiment fracture ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "Force 5: even the workers are splitting in two", 29, ACCENT, True)
txt(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(0.6),
    "2026 tech-workforce survey — anxiety tracks the same line", 17, GREY, True)
rows = [
    ("Designers", "63% overwhelmed · 61% tired", ACCENT, 0.90),
    ("Researchers", "36% fear job loss · 51% anxious", ACCENT, 0.72),
    ("Founders", "only 15% anxious", GREEN, 0.24),
]
y = Inches(2.3)
for who, what, col, frac in rows:
    txt(s, Inches(0.9), y, Inches(2.6), Inches(0.7), who, 20, INK, True, PP_ALIGN.LEFT,
        MSO_ANCHOR.MIDDLE)
    doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.6), y+Inches(0.05), Inches(8.0*frac), Inches(0.6),
           color=INK, fill=col, wdt=2)
    txt(s, Inches(3.75), y+Inches(0.03), Inches(8.2), Inches(0.65), what, 15, INK, True,
        PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE) if frac < 0.5 else \
    txt(s, Inches(3.75), y+Inches(0.03), Inches(8.0*frac-0.3), Inches(0.65), what, 15, PAPER, True,
        PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    y = Emu(int(y + Inches(1.1)))
txt(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(1.0),
    "The person who owns more judgment feels SAFER.\n"
    "Fear is now a position indicator — it tells you which track you're on.",
    19, ACCENT, True, PP_ALIGN.CENTER)

# ---------- Slide 9: the tool ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "\U0001F6E0 A tool to try this week", 32, ACCENT, True)
txt(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(0.6),
    "The Two-Track Sorter (10 minutes, one page)", 24, INK, True)
steps = [
    "1.  Name your role's core task in one verb-phrase (“write proposals”, “check forms”).",
    "2.  Ask the fork question: is AI PROFESSIONALISING my task, or DEMOCRATISING it?",
    "3.  If democratising → find the judgment call next to it a customer pays a human for.",
    "4.  If professionalising → double down; the 42%-faster raises are on this track.",
    "5.  Entry-level check: what senior judgment can you DEMONSTRATE to skip the deleted rung?",
]
y = Inches(2.2)
for st in steps:
    doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), y, Inches(11.5), Inches(0.72),
           color=BLUE, fill=RGBColor(0xF0,0xF5,0xFB), wdt=2)
    txt(s, Inches(1.15), y+Inches(0.06), Inches(11.0), Inches(0.62), st, 16, INK, False,
        PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    y = Emu(int(y + Inches(0.85)))
txt(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.8),
    "Ask your AI: “which parts of my task are you already better at, and which still need my judgment?”",
    16, ACCENT, True, PP_ALIGN.CENTER)

# ---------- Slide 10: closing / sources ----------
s = slide()
txt(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(1.0),
    "The takeaway", 40, ACCENT, True)
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.9), Inches(11.3), Inches(2.5),
       color=INK, fill=RGBColor(0xFF,0xF3,0xD9), wdt=3)
txt(s, Inches(1.4), Inches(2.1), Inches(10.5), Inches(2.1),
    "The AI-skills premium isn't a rising tide — it's a WEDGE.\n\n"
    "It lifts the judgment-heavy, customer-facing side and presses\n"
    "down on the commoditised, entry-level side.\n"
    "Know which side of every fork your work is on — and walk toward the judgment.",
    22, INK, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.9), Inches(4.7), Inches(11.5), Inches(2.2),
    "Sources: PwC 2026 Global AI Jobs Barometer · SUCCESS · Lightcast × America Succeeds "
    "(Durable by Design) ·\nMcKinsey · Lenny's Newsletter tech-workforce survey · ADP Research · "
    "Resume-Now · WEF · Forbes.\n(Full links in the accompanying README.)",
    15, GREY, False, PP_ALIGN.CENTER)

import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "Two_Track_Labour_Market_2026.pptx")
prs.save(out)
print("saved deck with", len(prs.slides._sldIdLst), "slides ->", out)
