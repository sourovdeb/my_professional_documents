#!/usr/bin/env python3
"""Doodle-style deck: AI & Durable Skills — 2026 trend report."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# palette
INK   = RGBColor(0x1D, 0x1D, 0x2B)   # near-black ink
PAPER = RGBColor(0xFB, 0xF7, 0xEE)   # warm paper
ACCENT= RGBColor(0xE8, 0x5D, 0x2E)   # marker orange
BLUE  = RGBColor(0x2E, 0x6B, 0xE8)   # marker blue
GREEN = RGBColor(0x2E, 0xA0, 0x5A)   # marker green
MUTE  = RGBColor(0x6B, 0x66, 0x5C)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def bg(slide, color=PAPER):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def box(slide, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font="Comic Sans MS", italic=False, line_spacing=1.05):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = color; f.name = font
    return tb

def shape(slide, kind, x, y, w, h, fill=None, line=INK, line_w=2.25):
    sp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def marker_underline(slide, x, y, w, color=ACCENT):
    ln = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 0.12, fill=color, line=color, line_w=0.5)
    return ln

# ---------- Slide 1: cover ----------
s = prs.slides.add_slide(BLANK); bg(s)
shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.28, fill=ACCENT, line=ACCENT, line_w=0.5)
box(s, 0.9, 1.6, 11.5, 1.6, "The Great Skill Split", size=60, bold=True, color=INK)
marker_underline(s, 0.95, 2.85, 6.2, ACCENT)
box(s, 0.95, 3.15, 11.4, 1.2,
    "Why your “human” skills just got a raise", size=30, color=BLUE, bold=True)
box(s, 0.95, 4.2, 11.4, 1.4,
    "AI is not flattening the value of skills — it is SPLITTING it.\n"
    "Routine work gets cheaper. Judgment, taste & steering AI get a pay rise.",
    size=20, color=INK)
box(s, 0.95, 6.4, 11.4, 0.6, "5-minute read  ·  Trend monitor  ·  July 2026",
    size=16, color=MUTE, italic=True)

# ---------- Slide 2: two doors doodle ----------
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0.7, 0.4, 12, 0.9, "One labour market, two doors", size=36, bold=True, color=INK)
marker_underline(s, 0.75, 1.25, 5.6, ACCENT)
box(s, 0.75, 1.5, 12, 0.6, "Same person, same skill. What decides your door? Whether you ADD judgment — or just hand the task over.",
    size=17, color=MUTE, italic=True)
# funnel node
shape(s, MSO_SHAPE.OVAL, 5.7, 2.25, 2.1, 0.9, fill=None, line=INK, line_w=2.5)
box(s, 5.7, 2.45, 2.1, 0.6, "YOU + a skill", size=15, bold=True, align=PP_ALIGN.CENTER)
# left door (routine)
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.3, 3.9, 4.6, 2.5, fill=None, line=MUTE, line_w=2.5)
box(s, 1.5, 4.05, 4.2, 0.7, "ROUTINE DOOR", size=22, bold=True, color=MUTE, align=PP_ALIGN.CENTER)
box(s, 1.5, 4.75, 4.2, 1.6,
    "“AI can do this for you”\n\nvalue drifts DOWN ↓\nwages flat, roles thin\n“prompt it and pray”",
    size=15, color=MUTE, align=PP_ALIGN.CENTER)
# right door (judgment)
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 7.4, 3.9, 4.6, 2.5, fill=None, line=GREEN, line_w=3)
box(s, 7.6, 4.05, 4.2, 0.7, "JUDGMENT DOOR", size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
box(s, 7.6, 4.75, 4.2, 1.6,
    "“AI needs a human on the loop”\n\nvalue climbs UP ↑\n+62% wage premium\n“prompt it, THEN judge it”",
    size=15, color=INK, align=PP_ALIGN.CENTER)
# arrows
a1 = shape(s, MSO_SHAPE.DOWN_ARROW, 3.35, 3.15, 0.5, 0.75, fill=MUTE, line=MUTE, line_w=1)
a2 = shape(s, MSO_SHAPE.DOWN_ARROW, 9.45, 3.15, 0.5, 0.75, fill=GREEN, line=GREEN, line_w=1)

# ---------- Slide 3: the numbers ----------
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0.7, 0.4, 12, 0.9, "The 2026 numbers that matter", size=36, bold=True, color=INK)
marker_underline(s, 0.75, 1.25, 5.9, ACCENT)
stats = [
    ("+62%", "AI-skill wage premium\n(was 57% → ~25%)", ACCENT),
    ("~8×", "AI jobs grow faster\n69% vs 9% market", BLUE),
    ("51%", "of AI jobs are now\nOUTSIDE IT", GREEN),
    ("7×", "entry-level roles want\nSENIOR skills", ACCENT),
    ("22%", "of workers feel their\njob is safe", BLUE),
    ("50%", "of orgs to run\n“AI-free” thinking tests", GREEN),
]
x0, y0, w, h, gx, gy = 0.85, 1.7, 3.7, 2.35, 0.35, 0.35
for i, (big, small, col) in enumerate(stats):
    cx = x0 + (i % 3) * (w + gx)
    cy = y0 + (i // 3) * (h + gy)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, w, h, fill=None, line=col, line_w=2.5)
    box(s, cx, cy + 0.18, w, 1.0, big, size=46, bold=True, color=col, align=PP_ALIGN.CENTER)
    box(s, cx, cy + 1.35, w, 0.9, small, size=15, color=INK, align=PP_ALIGN.CENTER)
box(s, 0.85, 6.75, 12, 0.5, "Sources: PwC 2026 AI Jobs Barometer · Lightcast · Gartner · ADP Research",
    size=12, color=MUTE, italic=True)

# ---------- Slide 4: the other half — skillidays ----------
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0.7, 0.4, 12.2, 0.9, "The human antidote: “skillidays”", size=36, bold=True, color=INK)
marker_underline(s, 0.75, 1.25, 6.4, GREEN)
box(s, 0.75, 1.45, 11.8, 0.9,
    "Learning for identity, not just jobs. 48% of Europeans (Gen Z: 57%) plan to learn a skill on holiday.\n"
    "In an AI world, we crave what AI can’t do FOR us — embodied, experiential, human.",
    size=17, color=INK, italic=True)
bars = [("Languages", 30, BLUE), ("Cookery", 28, ACCENT),
        ("Wellness (yoga/dance)", 25, GREEN), ("Traditional crafts", 24, INK)]
bx, by, maxw = 3.9, 2.9, 7.4
for i, (label, pct, col) in enumerate(bars):
    yy = by + i * 0.92
    box(s, 0.75, yy - 0.02, 3.0, 0.6, label, size=17, bold=True, color=INK, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, bx, yy, maxw * pct / 30.0, 0.55, fill=col, line=col, line_w=0.5)
    box(s, bx + maxw * pct / 30.0 + 0.15, yy - 0.02, 1.4, 0.6, f"{pct}%", size=18, bold=True, color=col, anchor=MSO_ANCHOR.MIDDLE)
box(s, 0.75, 6.75, 12, 0.5, "Source: Mastercard survey, 27,000 travellers / 28 European countries (Euronews, Jul 2026)",
    size=12, color=MUTE, italic=True)

# ---------- Slide 5: the tool ----------
s = prs.slides.add_slide(BLANK); bg(s)
shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.28, fill=BLUE, line=BLUE, line_w=0.5)
box(s, 0.7, 0.5, 12, 0.9, "\U0001F6E0  The tool: the AI × Judgment Stack", size=34, bold=True, color=INK)
marker_underline(s, 0.75, 1.4, 7.2, BLUE)
# equation
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.4, 1.85, 10.5, 1.0, fill=None, line=INK, line_w=2.5)
box(s, 1.4, 2.02, 10.5, 0.7, "VALUE  =  AI-LEVERAGE  ×  JUDGMENT  ×  DOMAIN",
    size=26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
rows = [
    ("AI-LEVERAGE", "Do I use AI to do this 5–10× faster?", BLUE),
    ("JUDGMENT", "Do I catch what AI gets wrong or can’t see?", ACCENT),
    ("DOMAIN", "Do I know this field deeper than the model?", GREEN),
]
for i, (k, q, col) in enumerate(rows):
    yy = 3.15 + i * 0.82
    box(s, 1.5, yy, 3.2, 0.6, k + "  [1–5]", size=18, bold=True, color=col)
    box(s, 5.0, yy, 7.0, 0.6, q, size=17, color=INK)
box(s, 1.5, 5.75, 10.5, 1.3,
    "Multiply the three (max 125):\n"
    "  •  < 20   Routine-Door risk — add a judgment or domain layer\n"
    "  •  20–60  Mixed — level up your lowest score\n"
    "  •  > 60   Judgment-Door moat — go deeper",
    size=16, color=INK)

# ---------- Slide 6: takeaway ----------
s = prs.slides.add_slide(BLANK); bg(s, INK)
box(s, 1.0, 1.5, 11.3, 1.2, "The move for 2026", size=40, bold=True, color=PAPER)
marker_underline(s, 1.05, 2.6, 4.5, ACCENT)
box(s, 1.0, 3.0, 11.3, 2.2,
    "Not AI OR human skills.\n\nAI  ×  critical thinking  ×  domain expertise.",
    size=30, bold=True, color=PAPER, line_spacing=1.15)
box(s, 1.0, 5.6, 11.3, 1.2,
    "Stack judgment on top of automation. Keep the embodied, human learning alive.\n"
    "The real meta-skill is learning how to learn — because the tools change every year.",
    size=18, color=RGBColor(0xE8, 0xE2, 0xD5), italic=True)

out = "/home/user/my_professional_documents/free_education/2026-07-18-durable-skills-vs-ai/AI_Durable_Skills_Trends_2026.pptx"
prs.save(out)
print("saved", out, len(prs.slides.__iter__.__self__._sldIdLst), "slides")
