#!/usr/bin/env python3
"""Doodle-style deck: You Just Got Promoted to 'Agent Boss' — 21 Jul 2026 trend brief."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# palette
INK   = RGBColor(0x1D, 0x1D, 0x2B)   # near-black ink
PAPER = RGBColor(0xFB, 0xF7, 0xEE)   # warm paper
ACCENT= RGBColor(0xE8, 0x5D, 0x2E)   # marker orange
BLUE  = RGBColor(0x2E, 0x6B, 0xE8)   # marker blue
GREEN = RGBColor(0x2E, 0xA0, 0x5A)   # marker green
MUTE  = RGBColor(0x6B, 0x66, 0x5C)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def bg(slide, color=PAPER):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def box(slide, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font="Comic Sans MS", italic=False, line_spacing=1.05):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = color; f.name = font
    return tb

def shape(slide, kind, x, y, w, h, fill=None, line=INK, line_w=2.25):
    sp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def marker_underline(slide, x, y, w, color=ACCENT):
    return shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 0.12, fill=color, line=color, line_w=0.5)

# ---------- Slide 1: cover ----------
s = prs.slides.add_slide(BLANK); bg(s)
shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.28, fill=ACCENT, line=ACCENT, line_w=0.5)
box(s, 0.9, 1.4, 11.5, 1.6, "You just got promoted", size=54, bold=True, color=INK)
box(s, 0.9, 2.5, 11.5, 1.0, "to “Agent Boss”", size=54, bold=True, color=ACCENT)
marker_underline(s, 0.95, 3.65, 6.8, BLUE)
box(s, 0.95, 3.95, 11.4, 1.4,
    "As AI agents do the execution, the #1 human skill flipped:\n"
    "not producing work — CHECKING the machine’s work.",
    size=22, color=INK)
box(s, 0.95, 6.4, 11.4, 0.6, "5-minute read  ·  Trend monitor  ·  21 July 2026",
    size=16, color=MUTE, italic=True)

# ---------- Slide 2: the work moved up one floor ----------
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0.7, 0.4, 12, 0.9, "The work moved up one floor", size=36, bold=True, color=INK)
marker_underline(s, 0.75, 1.25, 6.0, ACCENT)
box(s, 0.75, 1.45, 11.8, 0.6, "You used to sit on the bottom floor doing the task. The agent sits there now — you moved up.",
    size=17, color=MUTE, italic=True)
# top floor: agent boss
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 3.6, 2.3, 6.1, 1.7, fill=None, line=GREEN, line_w=3)
box(s, 3.7, 2.4, 5.9, 0.6, "YOU — the AGENT BOSS", size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
box(s, 3.7, 3.05, 5.9, 0.9, "aim the agent  ·  JUDGE the output ★  ·  own the outcome",
    size=16, color=INK, align=PP_ALIGN.CENTER)
box(s, 9.8, 2.55, 3.3, 0.6, "← the new #1 skill", size=16, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE)
# arrow down
shape(s, MSO_SHAPE.DOWN_ARROW, 6.35, 4.1, 0.55, 0.8, fill=MUTE, line=MUTE, line_w=1)
box(s, 7.0, 4.25, 5.0, 0.5, "delegates execution", size=14, color=MUTE, italic=True, anchor=MSO_ANCHOR.MIDDLE)
# bottom floor: agent
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 3.6, 5.05, 6.1, 1.5, fill=None, line=MUTE, line_w=2.5)
box(s, 3.7, 5.15, 5.9, 0.6, "AI AGENT — the doer", size=22, bold=True, color=MUTE, align=PP_ALIGN.CENTER)
box(s, 3.7, 5.8, 5.9, 0.6, "drafts  ·  builds  ·  fetches  ·  analyses", size=16, color=INK, align=PP_ALIGN.CENTER)

# ---------- Slide 3: the numbers ----------
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0.7, 0.4, 12, 0.9, "The numbers that shifted this week", size=36, bold=True, color=INK)
marker_underline(s, 0.75, 1.25, 6.6, ACCENT)
stats = [
    ("50%", "say quality-control of\nAI output is the #1 skill", ACCENT),
    ("46%", "critical thinking\n— the #2 skill", BLUE),
    ("62%", "AI-skill wage premium\n(“144% more demand”)", GREEN),
    ("+40%", "top per-skill pay:\nmachine learning", ACCENT),
    ("1.3M", "new AI job types\ncreated in 2 yrs (LinkedIn)", BLUE),
    ("49%", "of AI use is cognitive\nwork, not typing", GREEN),
]
x0, y0, w, h, gx, gy = 0.85, 1.7, 3.7, 2.35, 0.35, 0.35
for i, (big, small, col) in enumerate(stats):
    cx = x0 + (i % 3) * (w + gx)
    cy = y0 + (i // 3) * (h + gy)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, w, h, fill=None, line=col, line_w=2.5)
    box(s, cx, cy + 0.18, w, 1.0, big, size=44, bold=True, color=col, align=PP_ALIGN.CENTER)
    box(s, cx, cy + 1.35, w, 0.9, small, size=15, color=INK, align=PP_ALIGN.CENTER)
box(s, 0.85, 6.75, 12, 0.5, "Sources: Microsoft 2026 Work Trend Index · PwC AI Jobs Barometer · LinkedIn · index.dev",
    size=12, color=MUTE, italic=True)

# ---------- Slide 4: mood — both true at once ----------
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0.7, 0.4, 12.2, 0.9, "Energised AND uneasy — at once", size=36, bold=True, color=INK)
marker_underline(s, 0.75, 1.25, 6.6, BLUE)
box(s, 0.75, 1.45, 11.8, 0.6, "Both columns are true for the same workers this month (Microsoft WTI 2026).",
    size=17, color=MUTE, italic=True)
# upside col
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, 2.3, 5.5, 4.0, fill=None, line=GREEN, line_w=3)
box(s, 1.05, 2.45, 5.2, 0.6, "THE UPSIDE", size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
box(s, 1.15, 3.2, 5.0, 3.0,
    "66%  spend more time on\n        high-value work\n\n"
    "58%  produce work they\n        couldn’t make a year ago\n\n"
    "49%  of AI use is cognitive\n        (analysis, problem-solving)",
    size=17, color=INK)
# fear col
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 6.9, 2.3, 5.5, 4.0, fill=None, line=ACCENT, line_w=3)
box(s, 7.05, 2.45, 5.2, 0.6, "THE FEAR", size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
box(s, 7.15, 3.2, 5.0, 3.0,
    "65%  fear falling behind\n        without adapting to AI\n\n"
    "45%  say redesigning work\n        feels riskier than standing still\n\n"
    "22%  (ADP) still feel their\n        job is safe",
    size=17, color=INK)

# ---------- Slide 5: the tool ----------
s = prs.slides.add_slide(BLANK); bg(s)
shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.28, fill=BLUE, line=BLUE, line_w=0.5)
box(s, 0.7, 0.5, 12, 0.9, "\U0001F6E0  The tool: B.O.S.S.S. — inspect before you ship", size=30, bold=True, color=INK)
marker_underline(s, 0.75, 1.4, 8.6, BLUE)
box(s, 0.75, 1.55, 11.8, 0.6, "You can’t be an agent boss without a fast way to inspect. Five gates, 2 minutes, every output.",
    size=16, color=MUTE, italic=True)
rows = [
    ("B  BIAS", "Whose view is baked in? What’s missing?", ACCENT),
    ("O  ORIGIN", "Can I trace every number/claim to a source?", BLUE),
    ("S  SANITY", "Does it pass a back-of-envelope gut check?", GREEN),
    ("S  STAKES", "If this is wrong, who gets hurt — and how badly?", ACCENT),
    ("S  SECOND", "Ask the AI to argue the OPPOSITE, then judge.", BLUE),
]
for i, (k, q, col) in enumerate(rows):
    yy = 2.35 + i * 0.78
    box(s, 1.1, yy, 3.1, 0.6, k, size=19, bold=True, color=col, anchor=MSO_ANCHOR.MIDDLE)
    box(s, 4.4, yy, 8.0, 0.6, q, size=17, color=INK, anchor=MSO_ANCHOR.MIDDLE)
box(s, 1.1, 6.5, 11.4, 0.7, "Green on all five → ship.  Any red → you just earned your salary.",
    size=17, bold=True, color=GREEN, italic=True)

# ---------- Slide 6: takeaway ----------
s = prs.slides.add_slide(BLANK); bg(s, INK)
box(s, 1.0, 1.4, 11.3, 1.2, "The move for this week", size=40, bold=True, color=PAPER)
marker_underline(s, 1.05, 2.5, 5.2, ACCENT)
box(s, 1.0, 2.9, 11.3, 2.2,
    "The agent took the bottom floor.\n\nYou moved up to agent boss — rent paid in JUDGMENT.",
    size=28, bold=True, color=PAPER, line_spacing=1.15)
box(s, 1.0, 5.5, 11.3, 1.3,
    "Build the skill the market just named #1: validate what the machine produces.\n"
    "Keep one embodied human skill alive off the clock. Keep learning how to learn.",
    size=18, color=RGBColor(0xE8, 0xE2, 0xD5), italic=True)

out = "/home/user/my_professional_documents/free_education/2026-07-21-agent-boss-validation-skill/Agent_Boss_Trends_2026.pptx"
prs.save(out)
print("saved", out, len(prs.slides._sldIdLst), "slides")
