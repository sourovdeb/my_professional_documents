#!/usr/bin/env python3
"""Doodle-style deck: The Scoreboard Arrives — AI & Human Skills, 31 July 2026."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# palette (shared house style across the monitor series)
INK   = RGBColor(0x1D, 0x1D, 0x2B)   # near-black ink
PAPER = RGBColor(0xFB, 0xF7, 0xEE)   # warm paper
ACCENT= RGBColor(0xE8, 0x5D, 0x2E)   # marker orange
BLUE  = RGBColor(0x2E, 0x6B, 0xE8)   # marker blue
GREEN = RGBColor(0x2E, 0xA0, 0x5A)   # marker green
RED   = RGBColor(0xC0, 0x39, 0x2B)   # alarm red
MUTE  = RGBColor(0x6B, 0x66, 0x5C)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def bg(slide, color=PAPER):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def box(slide, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font="Comic Sans MS", italic=False, line_spacing=1.05):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = color; f.name = font
    return tb

def shape(slide, kind, x, y, w, h, fill=None, line=INK, line_w=2.25):
    sp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def underline(slide, x, y, w, color=ACCENT):
    return shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 0.12, fill=color, line=color, line_w=0.5)

# ---------- Slide 1: cover ----------
s = prs.slides.add_slide(BLANK); bg(s)
shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.28, fill=ACCENT, line=ACCENT, line_w=0.5)
box(s, 0.9, 1.5, 11.6, 1.6, "The Scoreboard Arrives", size=58, bold=True, color=INK)
underline(s, 0.95, 2.75, 6.6, ACCENT)
box(s, 0.95, 3.05, 11.5, 1.1,
    "The AI-jobs story just moved from forecast to measured", size=28, color=BLUE, bold=True)
box(s, 0.95, 4.15, 11.5, 1.5,
    "AI is now the #1 stated reason for US layoffs — four months running.\n"
    "The entry-level decline is no longer a prediction. It is a measured 13%.",
    size=20, color=INK)
box(s, 0.95, 6.4, 11.5, 0.6, "5-minute read  ·  Human-Skills trend monitor  ·  31 July 2026",
    size=16, color=MUTE, italic=True)

# ---------- Slide 2: soft signal -> hard scoreboard ----------
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0.7, 0.4, 12, 0.9, "Soft signal  →  hard scoreboard", size=36, bold=True, color=INK)
underline(s, 0.75, 1.25, 6.2, ACCENT)
box(s, 0.75, 1.5, 12, 0.6, "The evidence quality jumped: from what people FEAR to what the data SHOWS.",
    size=17, color=MUTE, italic=True)
# left card - before
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.9, 2.4, 5.4, 3.7, fill=None, line=MUTE, line_w=2.5)
box(s, 1.1, 2.55, 5.0, 0.7, "EARLIER 2026", size=20, bold=True, color=MUTE, align=PP_ALIGN.CENTER)
box(s, 1.1, 3.35, 5.0, 2.6,
    "“AI pays 62% more”      (survey)\n\n"
    "“AI skills +144% demand”  (job ads)\n\n"
    "“only 22% feel job-safe”  (sentiment)\n\n"
    "forecasts & feelings\n— easy to wave away",
    size=16, color=MUTE, align=PP_ALIGN.CENTER)
# arrow
shape(s, MSO_SHAPE.RIGHT_ARROW, 6.45, 3.9, 0.45, 0.7, fill=ACCENT, line=ACCENT, line_w=1)
# right card - now
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 7.05, 2.4, 5.4, 3.7, fill=None, line=RED, line_w=3)
box(s, 7.25, 2.55, 5.0, 0.7, "NOW (end of July)", size=20, bold=True, color=RED, align=PP_ALIGN.CENTER)
box(s, 7.25, 3.35, 5.0, 2.6,
    "“101,743 AI-cited layoffs YTD”   COUNTED\n\n"
    "“22–25s down 13% in exposed jobs”  MEASURED\n\n"
    "“AI = #1 layoff reason, 4 mo.”   RECORDED\n\n"
    "hard operational data\n— shows up in the ledger",
    size=16, color=INK, align=PP_ALIGN.CENTER)

# ---------- Slide 3: the layoff ledger ----------
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0.7, 0.4, 12, 0.9, "1 · The layoff ledger now names AI", size=34, bold=True, color=INK)
underline(s, 0.75, 1.25, 6.4, RED)
stats = [
    ("#1", "stated layoff reason\n4 months straight", RED),
    ("101,743", "AI-cited cuts\nin H1 2026", ACCENT),
    ("~2×", "all of 2025\n(was 54,836)", BLUE),
    ("31%", "of June's cuts\ntagged to AI", RED),
    ("+83%", "tech-sector cuts\nYoY (H1)", ACCENT),
    ("hiring", "Goldman: AI suppresses\nHIRING > firing", BLUE),
]
x0, y0, w, h, gx, gy = 0.85, 1.75, 3.7, 2.3, 0.35, 0.35
for i, (big, small, col) in enumerate(stats):
    cx = x0 + (i % 3) * (w + gx)
    cy = y0 + (i // 3) * (h + gy)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, w, h, fill=None, line=col, line_w=2.5)
    box(s, cx, cy + 0.2, w, 1.0, big, size=40, bold=True, color=col, align=PP_ALIGN.CENTER)
    box(s, cx, cy + 1.3, w, 0.9, small, size=15, color=INK, align=PP_ALIGN.CENTER)
box(s, 0.85, 6.8, 12, 0.5, "Source: Challenger, Gray & Christmas June 2026 report · Goldman Sachs",
    size=12, color=MUTE, italic=True)

# ---------- Slide 4: entry-level measured ----------
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0.7, 0.4, 12.2, 0.9, "2 · The entry-level squeeze is now measured", size=32, bold=True, color=INK)
underline(s, 0.75, 1.25, 6.8, RED)
box(s, 0.75, 1.45, 11.9, 0.8,
    "Stanford, on ADP payroll data for 25 million workers: employment change since late 2022.",
    size=17, color=INK, italic=True)
bars = [("Junior software devs (22–25)", -20, RED),
        ("All young workers (22–25)", -13, ACCENT),
        ("Older / less-exposed workers", 0, GREEN)]
bx, by, maxw = 5.4, 2.7, 6.4
for i, (label, pct, col) in enumerate(bars):
    yy = by + i * 1.05
    box(s, 0.75, yy - 0.02, 4.4, 0.6, label, size=16, bold=True, color=INK,
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    width = maxw * abs(pct) / 20.0 if pct != 0 else 0.25
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, bx, yy, max(width, 0.25), 0.6, fill=col, line=col, line_w=0.5)
    label_pct = f"{pct}%" if pct != 0 else "~flat"
    box(s, bx + max(width, 0.25) + 0.15, yy - 0.02, 1.6, 0.6, label_pct, size=18, bold=True,
        color=col, anchor=MSO_ANCHOR.MIDDLE)
box(s, 0.75, 6.05, 11.9, 0.9,
    "Declines cluster where AI AUTOMATES work — and are muted where it AUGMENTS. Dallas Fed agrees.",
    size=16, color=MUTE, italic=True)
box(s, 0.75, 6.85, 12, 0.4, "Source: Stanford (Brynjolfsson et al.) · Dallas Fed", size=12, color=MUTE, italic=True)

# ---------- Slide 5: the new gap ----------
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0.7, 0.4, 12, 0.9, "3 · Trained for today's AI — not tomorrow's job", size=30, bold=True, color=INK)
underline(s, 0.75, 1.2, 7.0, BLUE)
box(s, 0.75, 1.45, 11.8, 1.0,
    "The Conference Board (28 Jul): most organisations prepare workers for TODAY's AI —\n"
    "investing in AI literacy & current-role upskilling, but rarely in real reskilling.",
    size=18, color=INK)
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.4, 3.0, 10.5, 1.5, fill=None, line=ACCENT, line_w=3)
box(s, 1.6, 3.25, 10.1, 1.1,
    "So don't outsource your reskilling to your employer's training budget.\n"
    "Their horizon is this quarter's tool. Yours has to be the next role.",
    size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
box(s, 0.75, 5.0, 11.8, 1.6,
    "Durable-skill lists this week converge:\n"
    "   •  LinkedIn — AI literacy is #1, and means JUDGING the machine · soft skills take 7 of top 10\n"
    "   •  Resume-of-the-Future (O*NET) — reasoning, problem recognition, flexible thinking, judgment, focus\n"
    "   •  78% of ICT roles now embed AI skills → the differentiator is the human layer on top",
    size=16, color=INK)

# ---------- Slide 6: the tool ----------
s = prs.slides.add_slide(BLANK); bg(s)
shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.28, fill=BLUE, line=BLUE, line_w=0.5)
box(s, 0.7, 0.5, 12, 0.9, "\U0001F6E0  The tool: the Next-Role Reskill Ledger", size=32, bold=True, color=INK)
underline(s, 0.75, 1.4, 7.6, BLUE)
steps = [
    ("1  NAME", "your role as an AI agent would describe it in one line"),
    ("2  STRIP", "which part can AI already draft?  = your at-risk core"),
    ("3  LIFT", "the judgment the routine was feeding: validate · decide · persuade"),
    ("4  RESKILL", "1 durable skill + 1 NEXT-tool skill (not this quarter's tool)"),
    ("5  DATE IT", "90-day review — beat the skill half-life on purpose"),
]
for i, (k, q) in enumerate(steps):
    yy = 1.95 + i * 0.72
    box(s, 1.1, yy, 2.5, 0.6, k, size=19, bold=True, color=ACCENT)
    box(s, 3.7, yy, 8.6, 0.6, q, size=17, color=INK)
shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.1, 5.65, 11.2, 1.25, fill=None, line=GREEN, line_w=2.5)
box(s, 1.3, 5.8, 10.9, 1.0,
    "Sparring partner, not oracle — paste your conclusion back to any AI and prompt:\n"
    "“Argue the strongest case that I'm wrong. List the three assumptions I didn't check.”",
    size=16, color=INK, italic=True, align=PP_ALIGN.CENTER)

# ---------- Slide 7: takeaway ----------
s = prs.slides.add_slide(BLANK); bg(s, INK)
box(s, 1.0, 1.4, 11.3, 1.2, "The move for 31 July 2026", size=38, bold=True, color=PAPER)
underline(s, 1.05, 2.5, 5.0, ACCENT)
box(s, 1.0, 2.9, 11.3, 2.0,
    "AI literacy  ×  critical thinking  ×  domain  ×  reskill-ahead.",
    size=28, bold=True, color=PAPER, line_spacing=1.15)
box(s, 1.0, 4.5, 11.3, 2.0,
    "The scoreboard is real: AI now names itself in the layoff ledger, and the entry-level\n"
    "decline is measured, not feared. As routine work gets automated, value migrates to the\n"
    "judgment ceiling — and “AI literacy” now means judging the machine, not just running it.\n"
    "Reskill for the NEXT role. Keep the embodied, human learning alive. Learn how to learn.",
    size=17, color=RGBColor(0xE8, 0xE2, 0xD5), italic=True, line_spacing=1.15)

out = "/home/user/my_professional_documents/free_education/2026-07-31-ai-jobs-scoreboard/AI_Jobs_Scoreboard_2026.pptx"
prs.save(out)
print("saved", out, len(prs.slides._sldIdLst), "slides")
