#!/usr/bin/env python3
"""Build the 'Two-Track Labour Market' doodle deck (late-July 2026 trend brief)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# palette
INK    = RGBColor(0x22, 0x22, 0x2B)
PAPER  = RGBColor(0xFB, 0xF7, 0xEF)
ACCENT = RGBColor(0xE8, 0x55, 0x3A)  # warm orange-red
BLUE   = RGBColor(0x2F, 0x6F, 0xB0)
GREEN  = RGBColor(0x3B, 0x8C, 0x5A)
GREY   = RGBColor(0x8A, 0x86, 0x7E)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def txt(s, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font="Comic Sans MS", line_spacing=1.05):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = font
    return tb


def doodle(s, shape, x, y, w, h, color=INK, fill=None, wdt=2.5):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = color; sp.line.width = Pt(wdt)
    sp.shadow.inherit = False
    return sp


# ---------- Slide 1: title ----------
s = slide()
txt(s, Inches(0.9), Inches(1.4), Inches(11.6), Inches(1.6),
    "The Two-Track\nLabour Market", 54, INK, True, PP_ALIGN.LEFT)
doodle(s, MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(3.55), Inches(7.0), Inches(0.16),
       color=ACCENT, fill=ACCENT, wdt=1)
txt(s, Inches(0.95), Inches(3.8), Inches(11.6), Inches(1.0),
    "AI stopped being a pay bump. It's now splitting the workforce.", 26, ACCENT, True)
txt(s, Inches(0.95), Inches(4.7), Inches(11.6), Inches(1.0),
    "Late-July 2026 trend brief  ·  5-minute read  ·  27 July 2026", 20, GREY)
doodle(s, MSO_SHAPE.LEFT_RIGHT_ARROW, Inches(10.2), Inches(1.3), Inches(2.3), Inches(1.2),
       color=BLUE, fill=RGBColor(0xE6, 0xEF, 0xF7), wdt=3)
txt(s, Inches(10.2), Inches(1.55), Inches(2.3), Inches(0.7), "A  |  B", 24, BLUE,
    True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

# ---------- Slide 2: the one line ----------
s = slide()
txt(s, Inches(0.9), Inches(0.6), Inches(11.5), Inches(0.8), "The one line", 34, ACCENT, True)
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.7), Inches(11.3), Inches(3.6),
       color=INK, fill=RGBColor(0xFF, 0xF3, 0xD9), wdt=3)
txt(s, Inches(1.5), Inches(2.0), Inches(10.3), Inches(3.1),
    "Eight days ago the story was the WAGE PREMIUM (+62%).\n\n"
    "The fresher data flips it to VOLUME + STRUCTURE:\n"
    "AI-skill demand grew +144% in a year (market: +7%),\n"
    "and the market is forking into two tracks.", 26, INK, True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
txt(s, Inches(1.0), Inches(5.6), Inches(11.3), Inches(1.2),
    "The twist: the fork hits EARLIEST at the bottom — AI is eating the junior on-ramp.",
    22, ACCENT, True, PP_ALIGN.CENTER)

# ---------- Slide 3: what changed since 19 July ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.6), Inches(0.8),
    "What changed since the 19 July brief", 32, ACCENT, True)
rows = [
    ("Wage premium (PwC)", "+62% (from 57%)", "+62% confirmed", "stable", GREY),
    ("AI-skill DEMAND growth", "“outside IT”", "+144% YoY vs +7%", "NEW LEAD", ACCENT),
    ("Premium dispersion", "one number", "118% → 16% by sector", "a SPLIT", ACCENT),
    ("Critical-thinking demand", "#1 skill (73%)", "triple-digit YoY growth", "now a rate", BLUE),
    ("Entry-level roles", "not tracked", "7× need senior judgment", "on-ramp", ACCENT),
    ("Top-tier satisfaction", "22% job-safe", "49.7% → 41.7% in 1 qtr", "cracking", GREY),
    ("Skilliday adoption", "57% plan to", "37% already BOOKED", "intent→action", GREEN),
]
y = Inches(1.5)
for label, base, now, tag, col in rows:
    doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), y, Inches(11.5), Inches(0.68),
           color=col, fill=PAPER, wdt=1.5)
    txt(s, Inches(1.05), y + Inches(0.05), Inches(3.5), Inches(0.6), label, 14, INK, True,
        PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    txt(s, Inches(4.5), y + Inches(0.05), Inches(2.6), Inches(0.6), base, 13, GREY,
        False, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    txt(s, Inches(7.1), y + Inches(0.05), Inches(3.4), Inches(0.6), now, 13, INK,
        True, PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    txt(s, Inches(10.5), y + Inches(0.05), Inches(1.8), Inches(0.6), tag, 13, col,
        True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    y = Emu(int(y + Inches(0.79)))

# ---------- Slide 4: the two tracks ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "One market becoming two", 32, ACCENT, True)
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.7), Inches(5.4), Inches(4.5),
       color=BLUE, fill=RGBColor(0xE6, 0xEF, 0xF7), wdt=4)
txt(s, Inches(1.2), Inches(1.9), Inches(4.8), Inches(0.7), "TRACK A", 26, BLUE, True)
txt(s, Inches(1.2), Inches(2.6), Inches(4.9), Inches(3.4),
    "“AI × judgment”\n\n• demand +144% YoY\n• premium up to +118%\n"
    "• judgment demanded\n  EARLY in a career\n• you validate & decide", 19, INK)
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.7), Inches(5.4), Inches(4.5),
       color=GREY, fill=RGBColor(0xEF, 0xEC, 0xE4), wdt=3)
txt(s, Inches(7.3), Inches(1.9), Inches(4.8), Inches(0.7), "TRACK B", 26, GREY, True)
txt(s, Inches(7.3), Inches(2.6), Inches(4.9), Inches(3.4),
    "“routine task”\n\n• demand flat / ↓\n• premium ~16% or none\n"
    "• the tasks AI now\n  drafts on its own\n• highest displacement risk", 19, INK)
txt(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.8),
    "The gap between them is exactly what “learning how to learn” closes.",
    20, ACCENT, True, PP_ALIGN.CENTER)

# ---------- Slide 5: demand chart ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "The new lead: demand, not just pay", 32, ACCENT, True)
base_y = Inches(6.0)
def bar(x, h_in, label, val, color):
    doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Emu(int(base_y - Inches(h_in))),
           Inches(1.6), Inches(h_in), color=INK, fill=color, wdt=2)
    txt(s, x - Inches(0.3), Emu(int(base_y - Inches(h_in) - Inches(0.55))), Inches(2.2),
        Inches(0.5), val, 22, INK, True, PP_ALIGN.CENTER)
    txt(s, x - Inches(0.3), base_y + Inches(0.05), Inches(2.2), Inches(0.8), label, 14, GREY,
        False, PP_ALIGN.CENTER)
bar(Inches(1.6), 0.35, "Overall market", "+7%", GREY)
bar(Inches(5.0), 3.6, "AI-skill jobs", "+144%", ACCENT)
bar(Inches(8.4), 2.8, "Critical thinking", "triple-\ndigit", BLUE)
txt(s, Inches(10.6), Inches(1.6), Inches(2.5), Inches(3.5),
    "AI-skill jobs\ngrowing ~8×\nfaster than\nthe market.\n\n(PwC 2026,\n1bn+ postings)", 15, INK)

# ---------- Slide 6: the junior on-ramp ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "AI is eating the junior on-ramp", 32, ACCENT, True)
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.7), Inches(5.4), Inches(4.4),
       color=GREY, fill=RGBColor(0xEF, 0xEC, 0xE4), wdt=2.5)
txt(s, Inches(1.2), Inches(1.9), Inches(4.8), Inches(0.6), "2022 JUNIOR", 22, GREY, True)
txt(s, Inches(1.2), Inches(2.6), Inches(4.8), Inches(3.3),
    "Months reconciling\nspreadsheets.\n\nDull — but it taught\nwhat “wrong data”\nfeels like.\n\n"
    "The apprenticeship.", 20, INK)
doodle(s, MSO_SHAPE.RIGHT_ARROW, Inches(6.45), Inches(3.4), Inches(0.9), Inches(0.9),
       color=ACCENT, fill=ACCENT, wdt=1)
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.7), Inches(5.0), Inches(4.4),
       color=BLUE, fill=RGBColor(0xE6, 0xEF, 0xF7), wdt=3)
txt(s, Inches(7.8), Inches(1.9), Inches(4.4), Inches(0.6), "2026 JUNIOR", 22, BLUE, True)
txt(s, Inches(7.8), Inches(2.6), Inches(4.4), Inches(3.3),
    "Agent reconciles it\nbefore they arrive.\n\nDay one: validate the\noutput & brief a\ndirector.\n\n"
    "7× more likely to need\nsenior judgment.", 20, INK)
txt(s, Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.8),
    "The on-ramp didn't vanish — it moved up the wall. Build the judgment layer on purpose.",
    19, ACCENT, True, PP_ALIGN.CENTER)

# ---------- Slide 7: the mood ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "The mood: satisfaction starts to crack", 30, ACCENT, True)
stats = [
    ("41.7%", "“extremely satisfied”\n— down from 49.7%\nin one quarter", ACCENT),
    ("22%", "still strongly agree\ntheir job is safe\n(ADP, unchanged)", GREY),
    ("30.8%", "of EXECUTIVES most\nAI-anxious — 2.5×\nany other level", BLUE),
    ("53%", "of AI mentions in\nreviews now negative\n(was 55% positive)", ACCENT),
]
x = Inches(0.9)
for val, body, col in stats:
    doodle(s, MSO_SHAPE.OVAL, x, Inches(1.9), Inches(2.7), Inches(2.7),
           color=col, fill=PAPER, wdt=3)
    txt(s, x, Inches(2.5), Inches(2.7), Inches(0.9), val, 32, col, True, PP_ALIGN.CENTER)
    txt(s, x, Inches(3.35), Inches(2.7), Inches(1.1), body, 12, INK, False, PP_ALIGN.CENTER)
    x = Emu(int(x + Inches(3.0)))
txt(s, Inches(0.9), Inches(5.1), Inches(11.5), Inches(1.6),
    "Economy (42%) is still cited as a bigger threat than AI (17%) — but AI is climbing.\n\n"
    "Fear is a bad strategy, a useful signal: move from doing the task to owning the judgment.",
    17, INK)

# ---------- Slide 8: skilliday + the tool ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "The counter-current + your move", 30, GREEN, True)
doodle(s, MSO_SHAPE.SUN, Inches(10.7), Inches(0.4), Inches(1.9), Inches(1.9),
       color=RGBColor(0xE0, 0xA5, 0x2A), fill=RGBColor(0xFF, 0xE7, 0x9c), wdt=3)
txt(s, Inches(0.9), Inches(1.35), Inches(6.0), Inches(0.6),
    "“Skilliday”: intent → action", 20, GREEN, True)
txt(s, Inches(0.9), Inches(2.0), Inches(5.6), Inches(3.0),
    "• ~half of Europeans plan;\n  37% ALREADY BOOKED\n• 82% say it's more memorable\n"
    "• Languages 30% · cooking 28%\n  · winemaking/cheese 28%\n\n"
    "Same instinct as work:\nreach for something embodied\n& hard to automate.", 16, INK)
txt(s, Inches(6.9), Inches(1.35), Inches(5.6), Inches(0.6),
    "\U0001f6e0 The Two-Track Self-Audit", 20, ACCENT, True)
steps = [
    "1. List the 5 tasks that fill your week.",
    "2. Sort each: Track A (judgment) or B (routine).",
    "3. Count Track B — your automation exposure.",
    "4. For each B, name the A-judgment it frees.",
    "5. Grow 1 durable + 1 technical. Review in 90 days.",
]
y = Inches(2.0)
for st in steps:
    doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), y, Inches(5.6), Inches(0.62),
           color=BLUE, fill=RGBColor(0xF0, 0xF5, 0xFB), wdt=1.5)
    txt(s, Inches(7.05), y + Inches(0.04), Inches(5.3), Inches(0.54), st, 13, INK, False,
        PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    y = Emu(int(y + Inches(0.72)))
txt(s, Inches(6.9), Inches(5.7), Inches(5.6), Inches(1.2),
    "Free tool: use any AI as a sparring partner —\n“Argue I'm wrong; list 3 assumptions I skipped.”",
    14, ACCENT, True)

# ---------- Slide 9: takeaway + sources ----------
s = slide()
txt(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(1.0), "The takeaway", 40, ACCENT, True)
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.9), Inches(11.3), Inches(2.3),
       color=INK, fill=RGBColor(0xFF, 0xF3, 0xD9), wdt=3)
txt(s, Inches(1.4), Inches(2.1), Inches(10.5), Inches(1.9),
    "Don't compete on the routine layer AI already owns.\n"
    "Build the judgment layer — deliberately, and early.", 26, INK, True,
    PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.9), Inches(4.5), Inches(11.5), Inches(2.2),
    "Sources: PwC 2026 Global AI Jobs Barometer · SUCCESS (62% pay / 144% demand) ·\n"
    "Job Skills Report 2026 · Critical Core Skills 2026 · 4 Corner Resources Q2 2026 · ADP ·\n"
    "Mastercard / Euronews skilliday survey · WEF Future of Jobs 2026.  (Full links in the README.)",
    15, GREY, False, PP_ALIGN.CENTER)

OUT = "/home/user/my_professional_documents/weekly-briefings/2026-07-27-ai-two-track-labour-market/Two_Track_Labour_Market_2026.pptx"
prs.save(OUT)
print("saved deck with", len(prs.slides._sldIdLst), "slides ->", OUT)
