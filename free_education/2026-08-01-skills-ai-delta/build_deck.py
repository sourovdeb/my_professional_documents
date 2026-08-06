#!/usr/bin/env python3
"""Doodle-style deck: AI, Skills & Learning — 1 Aug 2026 quiet-day check-in.

Honest low-news edition: the 31 July 'Scoreboard' brief still stands; this
logs minor previously-unrecorded texture and carries forward the tool.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# palette (matches prior decks)
INK    = RGBColor(0x1D, 0x1D, 0x2B)
PAPER  = RGBColor(0xFB, 0xF7, 0xEE)
ACCENT = RGBColor(0xE8, 0x5D, 0x2E)
BLUE   = RGBColor(0x2E, 0x6B, 0xE8)
GREEN  = RGBColor(0x2E, 0xA0, 0x5A)
MUTE   = RGBColor(0x6B, 0x66, 0x5C)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def bg(slide, color=PAPER):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font="Comic Sans MS", italic=False, line_spacing=1.05):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = color; f.name = font
    return tb


def shape(slide, kind, x, y, w, h, fill=None, line=INK, line_w=2.25):
    sp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def underline(slide, x, y, w, color=ACCENT):
    return shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 0.12, fill=color, line=color, line_w=0.5)


# ---------- Slide 1: cover ----------
s = prs.slides.add_slide(BLANK); bg(s)
shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.28, fill=GREEN, line=GREEN, line_w=0.5)
box(s, 0.9, 1.5, 11.5, 1.5, "Quiet Day", size=60, bold=True, color=INK)
underline(s, 0.95, 2.75, 5.0, GREEN)
box(s, 0.95, 3.05, 11.4, 1.0, "The Scoreboard still stands", size=30, color=BLUE, bold=True)
box(s, 0.95, 4.1, 11.4, 1.4,
    "One day on from the 31 July brief, no significant new data landed.\n"
    "The big move was last week: forecast → measured. Today = honest confirmation.",
    size=20, color=INK)
box(s, 0.95, 6.4, 11.4, 0.6, "5-minute check-in  ·  Trend monitor  ·  1 August 2026",
    size=16, color=MUTE, italic=True)

# ---------- Slide 2: where the story stands (still current) ----------
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0.7, 0.4, 12.2, 0.9, "Where the story stands (31 Jul — still current)", size=30, bold=True, color=INK)
underline(s, 0.75, 1.25, 7.4, GREEN)
box(s, 0.75, 1.45, 11.9, 0.55, "The measured picture from last week hasn't been superseded — it remains the plan:",
    size=17, color=MUTE, italic=True)
rows = [
    ("AI = #1 stated layoff reason", "4 months running · ~101,743 AI-cited cuts H1 (≈2× 2025)", ACCENT),
    ("Entry-level −13% (measured)", "22–25s in AI-exposed jobs; −20% junior devs — Stanford payroll", ACCENT),
    ("AI literacy = #1 rising skill", "soft skills take 7 of top 10 — LinkedIn 2026", BLUE),
    ("Trained for today's AI", "not tomorrow's job — Conference Board (28 Jul)", BLUE),
    ("+62% premium · 22% feel safe", "48% plan a skilliday — PwC · ADP · Mastercard", GREEN),
]
for i, (k, v, col) in enumerate(rows):
    yy = 2.25 + i * 0.85
    shape(s, MSO_SHAPE.OVAL, 0.8, yy + 0.16, 0.22, 0.22, fill=col, line=col, line_w=0.5)
    box(s, 1.2, yy, 5.0, 0.7, k, size=18, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    box(s, 6.2, yy, 6.5, 0.7, v, size=15, color=MUTE, anchor=MSO_ANCHOR.MIDDLE)
box(s, 0.75, 6.7, 12, 0.5, "✅  If you read the 31 July Scoreboard, you are up to date. Nothing today changes the plan.",
    size=14, color=GREEN, italic=True)

# ---------- Slide 3: what moved in 24h — the ladder ----------
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0.7, 0.4, 12, 0.9, "What moved in 24h: minor texture only", size=32, bold=True, color=INK)
underline(s, 0.75, 1.25, 6.3, ACCENT)
box(s, 0.75, 1.45, 11.9, 0.55, "No new flagship report. One detail worth logging: the \"22% feel safe\" number is a ladder.",
    size=16, color=MUTE, italic=True)
rungs = [("C-suite", 35, GREEN), ("Upper managers", 31, GREEN),
         ("Middle managers", 23, BLUE), ("Managers", 21, ACCENT), ("Rank-and-file", 18, ACCENT)]
bx, by, maxw = 4.3, 2.3, 6.6
for i, (lab, pct, col) in enumerate(rungs):
    yy = by + i * 0.72
    box(s, 0.6, yy - 0.02, 3.5, 0.55, lab, size=16, bold=True, color=INK,
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, bx, yy, maxw * pct / 35.0, 0.5, fill=col, line=col, line_w=0.5)
    box(s, bx + maxw * pct / 35.0 + 0.15, yy - 0.02, 1.4, 0.55, f"{pct}%", size=17, bold=True,
        color=col, anchor=MSO_ANCHOR.MIDDLE)
box(s, 0.6, 6.05, 12, 1.1,
    "Not new data — a finer cut of the same ADP survey. But it's the year in one shape:\n"
    "anxiety concentrates where the work is most automatable — the same rung Stanford's −13%\n"
    "payroll data shows actually shrinking. (Also, by role: researchers 51% anxious vs founders 15%.)",
    size=14, color=MUTE, italic=True)

# ---------- Slide 4: the tool ----------
s = prs.slides.add_slide(BLANK); bg(s)
shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.28, fill=BLUE, line=BLUE, line_w=0.5)
box(s, 0.7, 0.5, 12, 0.9, "\U0001F6E0  The tool: the Judgment Journal", size=34, bold=True, color=INK)
underline(s, 0.75, 1.4, 6.8, BLUE)
box(s, 0.75, 1.6, 11.9, 0.7,
    "Pairs with last week's Reskill Ledger. Gartner says half of orgs will soon test if you can think WITHOUT AI.",
    size=17, color=INK, italic=True)
rows = [
    ("1 · CAUGHT IT", "One thing AI got wrong this week that I caught & fixed.", ACCENT),
    ("2 · OVERRODE", "One time I chose against the AI's suggestion — and why.", BLUE),
    ("3 · NO-AI WIN", "One decision I made with zero AI help. My reasoning?", GREEN),
]
for i, (k, q, col) in enumerate(rows):
    yy = 2.7 + i * 0.95
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.2, yy, 3.1, 0.7, fill=None, line=col, line_w=2)
    box(s, 1.2, yy + 0.06, 3.1, 0.6, k, size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
    box(s, 4.6, yy, 7.5, 0.7, q, size=17, color=INK, anchor=MSO_ANCHOR.MIDDLE)
box(s, 1.2, 5.75, 10.9, 1.2,
    "5 minutes a week. After 8 weeks you hold a portfolio of judgment — concrete proof for the\n"
    "\"AI-free assessment\" era, and a mirror that keeps your critical thinking sharp.",
    size=16, color=INK)

# ---------- Slide 5: takeaway ----------
s = prs.slides.add_slide(BLANK); bg(s, INK)
box(s, 1.0, 1.4, 11.3, 1.2, "Bottom line", size=40, bold=True, color=PAPER)
underline(s, 1.05, 2.5, 3.4, GREEN)
box(s, 1.0, 2.95, 11.3, 1.8,
    "A genuinely quiet day.\n\nThe 31 July Scoreboard is the current state — no correction needed.",
    size=28, bold=True, color=PAPER, line_spacing=1.15)
box(s, 1.0, 5.5, 11.3, 1.3,
    "The move is unchanged: AI literacy × critical thinking × domain × reskill-ahead-of-the-tool.\n"
    "One fresh personal step: start a Judgment Journal alongside your Reskill Ledger.",
    size=18, color=RGBColor(0xE8, 0xE2, 0xD5), italic=True)

out = "/home/user/my_professional_documents/free_education/2026-08-01-skills-ai-delta/AI_Skills_Learning_Delta_2026-08.pptx"
prs.save(out)
print("saved", out, len(prs.slides._sldIdLst), "slides")
