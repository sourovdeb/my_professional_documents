#!/usr/bin/env python3
"""Build the 'Cognitive Debt — AI Deskilling' doodle deck (25 July 2026)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# palette
INK   = RGBColor(0x22, 0x22, 0x2B)
PAPER = RGBColor(0xFB, 0xF7, 0xEF)
ACCENT= RGBColor(0xE8, 0x55, 0x3A)  # warm orange-red
BLUE  = RGBColor(0x2F, 0x6F, 0xB0)
GREEN = RGBColor(0x3B, 0x8C, 0x5A)
GREY  = RGBColor(0x8A, 0x86, 0x7E)
RED   = RGBColor(0xC2, 0x3B, 0x2E)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height

def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid(); bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s

def txt(s, x, y, w, h, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font="Comic Sans MS", line_spacing=1.05):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.name = font
    return tb

def doodle(s, shape, x, y, w, h, color=INK, fill=None, wdt=2.5):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = color; sp.line.width = Pt(wdt)
    sp.shadow.inherit = False
    return sp

def blob_label(s, x, y, w, h, num, ncolor):
    doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, color=ncolor, fill=PAPER, wdt=3)
    doodle(s, MSO_SHAPE.OVAL, x-Inches(0.28), y-Inches(0.28), Inches(0.85), Inches(0.85),
           color=INK, fill=ncolor, wdt=2.5)
    txt(s, x-Inches(0.28), y-Inches(0.30), Inches(0.85), Inches(0.85), num, 26, PAPER,
        True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

# ---------- Slide 1: title ----------
s = slide()
txt(s, Inches(0.9), Inches(1.35), Inches(11.5), Inches(1.6),
    "Cognitive Debt", 62, INK, True, PP_ALIGN.LEFT)
doodle(s, MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(2.9), Inches(7.2), Inches(0.16),
       color=ACCENT, fill=ACCENT, wdt=1)
txt(s, Inches(0.95), Inches(3.15), Inches(11.6), Inches(1.2),
    "The tool that saves you time can erode the skill it saves.", 27, ACCENT, True)
txt(s, Inches(0.95), Inches(4.25), Inches(11.4), Inches(1.0),
    "AI-deskilling trend brief  ·  5-minute read  ·  25 July 2026", 20, GREY)
doodle(s, MSO_SHAPE.LIGHTNING_BOLT, Inches(10.9), Inches(1.0), Inches(1.4), Inches(2.0),
       color=ACCENT, fill=RGBColor(0xFF,0xE7,0x9c), wdt=3)
doodle(s, MSO_SHAPE.BLOCK_ARC, Inches(9.6), Inches(4.5), Inches(1.8), Inches(1.8),
       color=BLUE, fill=None, wdt=3)
txt(s, Inches(9.6), Inches(5.0), Inches(1.8), Inches(0.8), "use AI\nON PURPOSE", 15, BLUE,
    True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

# ---------- Slide 2: the one line ----------
s = slide()
txt(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.8), "The one line", 34, ACCENT, True)
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.6), Inches(11.3), Inches(3.5),
       color=INK, fill=RGBColor(0xFF,0xF3,0xD9), wdt=3)
txt(s, Inches(1.5), Inches(1.85), Inches(10.3), Inches(3.0),
    "Human skills don't just get MORE valuable\nas AI rises —\n\n"
    "they ATROPHY when you let AI do your\nthinking for you.", 30, INK, True,
    PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
txt(s, Inches(1.0), Inches(5.5), Inches(11.3), Inches(1.2),
    "New meta-move: not “use more AI.”  It's  use AI on purpose + keep your own reps.",
    24, ACCENT, True, PP_ALIGN.CENTER)

# ---------- Slide 3: old story vs new story ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "The plot twist since last month", 32, ACCENT, True)
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.7), Inches(5.4), Inches(4.3),
       color=GREY, fill=RGBColor(0xEF,0xEC,0xE4), wdt=2.5)
txt(s, Inches(1.2), Inches(1.9), Inches(4.8), Inches(0.6), "OLD STORY — automation", 20, GREY, True)
txt(s, Inches(1.2), Inches(2.7), Inches(4.8), Inches(3.1),
    "AI does the TASK.\n\nYou keep the skill,\nyou lose the busywork.\n\nvalue ▲", 22, INK)
doodle(s, MSO_SHAPE.RIGHT_ARROW, Inches(6.45), Inches(3.4), Inches(0.9), Inches(0.9),
       color=ACCENT, fill=ACCENT, wdt=1)
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(1.7), Inches(5.0), Inches(4.3),
       color=RED, fill=RGBColor(0xF9,0xE4,0xE0), wdt=3)
txt(s, Inches(7.8), Inches(1.9), Inches(4.4), Inches(0.6), "NEW STORY — cognitive debt", 20, RED, True)
txt(s, Inches(7.8), Inches(2.7), Inches(4.4), Inches(3.1),
    "AI does the THINKING.\n\nYou skip the reps,\nthe skill quietly decays.\n\n"
    "value ▼ (if unmanaged)", 22, INK)
txt(s, Inches(0.9), Inches(6.25), Inches(11.5), Inches(0.8),
    "Same tool. Opposite outcome — depending on how you use it.", 20, ACCENT, True, PP_ALIGN.CENTER)

# ---------- Slide 4: the evidence ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "The evidence landed this month", 32, ACCENT, True)
cards = [
    (Inches(0.9),  "Gartner", "50% of orgs will require\n“AI-FREE” skills\nassessments thru 2026", RED),
    (Inches(4.05), "Anthropic RCT", "AI group 50% vs\nhand-coding 67%\nrecall (“never-skilling”)", BLUE),
    (Inches(7.2),  "Wharton", "“cognitive surrender”:\nadopt the WRONG AI\nanswer >80% of the time", ACCENT),
    (Inches(10.35),"Radiology", "unaided tumor detection\n−6% in 3 months;\n74% clinicians worried", GREEN),
]
for x, head, body, col in cards:
    doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.9), Inches(2.75), Inches(3.6),
           color=col, fill=PAPER, wdt=3)
    txt(s, x+Inches(0.15), Inches(2.15), Inches(2.45), Inches(0.7), head, 19, col, True, PP_ALIGN.CENTER)
    txt(s, x+Inches(0.2), Inches(3.0), Inches(2.35), Inches(2.4), body, 15, INK, False, PP_ALIGN.CENTER)
txt(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(1.0),
    "BCG names it: the skills leaders rate most critical — judgment, problem-framing,\n"
    "creative thinking — are the ones most at risk of erosion.", 17, INK, True, PP_ALIGN.CENTER)

# ---------- Slide 5: demand got louder ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "Meanwhile demand got LOUDER, not quieter", 30, ACCENT, True)
base_y = Inches(5.0)
def bar(x, h_in, label, val, color):
    doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Emu(int(base_y - Inches(h_in))),
           Inches(1.4), Inches(h_in), color=INK, fill=color, wdt=2)
    txt(s, x-Inches(0.2), Emu(int(base_y - Inches(h_in) - Inches(0.55))), Inches(1.8),
        Inches(0.5), val, 20, INK, True, PP_ALIGN.CENTER)
    txt(s, x-Inches(0.2), base_y+Inches(0.05), Inches(1.8), Inches(0.9), label, 13, GREY, False, PP_ALIGN.CENTER)
bar(Inches(1.3), 0.5, "all roles", "+6%", GREY)
bar(Inches(3.6), 3.2, "AI-skill roles", "+134%", ACCENT)
txt(s, Inches(6.0), Inches(1.7), Inches(6.6), Inches(3.6),
    "• AI-skill jobs growing ~8× faster\n  than the market (69% vs 9%).\n\n"
    "• WEF: ~170M roles created, ~92M\n  displaced by 2030; half the\n  workforce needs reskilling in 4 yrs.\n\n"
    "• PwC wage premium holds at 62%\n  (was 57%).", 18, INK)
txt(s, Inches(0.9), Inches(6.2), Inches(11.5), Inches(0.9),
    "The premium isn't for the word “AI” on your résumé — it's for driving the tool AND catching it when it's wrong.",
    16, BLUE, True, PP_ALIGN.CENTER)

# ---------- Slide 6: the paradox ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "The paradox: active users are MORE engaged", 30, GREEN, True)
for x, val, body, col in [
    (Inches(1.1), "30%", "of DAILY AI users are\nfully engaged at work", GREEN),
    (Inches(4.6), "14%", "of those who NEVER\nuse AI are", GREY),
    (Inches(8.1), "44%", "“thriving at work” —\ndown from 66% in 2024", RED),
]:
    doodle(s, MSO_SHAPE.OVAL, x, Inches(1.9), Inches(3.0), Inches(3.0), color=col, fill=PAPER, wdt=3)
    txt(s, x, Inches(2.6), Inches(3.0), Inches(1.0), val, 44, col, True, PP_ALIGN.CENTER)
    txt(s, x, Inches(3.75), Inches(3.0), Inches(1.1), body, 14, INK, False, PP_ALIGN.CENTER)
txt(s, Inches(0.9), Inches(5.4), Inches(11.5), Inches(1.6),
    "Only 22% feel their job is safe; 99% of CEOs expect AI layoffs within 2 years.\n\n"
    "Read together: avoidance loses, and surrender loses. Engaged OWNERSHIP wins.",
    19, INK, True, PP_ALIGN.CENTER)

# ---------- Slide 7: the updated formula ----------
s = slide()
txt(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8),
    "The updated formula: add a denominator", 32, ACCENT, True)
boxes = [("AI\nleverage", BLUE), ("CRITICAL\nTHINKING", ACCENT), ("DOMAIN\nEXPERTISE", GREEN)]
x = Inches(1.0)
for i, (label, col) in enumerate(boxes):
    doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.9), Inches(2.7), Inches(1.9),
           color=col, fill=PAPER, wdt=4)
    txt(s, x, Inches(2.3), Inches(2.7), Inches(1.1), label, 22, col, True,
        PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    if i < 2:
        txt(s, Emu(int(x+Inches(2.7))), Inches(2.1), Inches(0.85), Inches(1.5), "×", 48, INK,
            True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    x = Emu(int(x + Inches(3.55)))
doodle(s, MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.15), Inches(11.3), Inches(0.05),
       color=INK, fill=INK, wdt=1)
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.0), Inches(4.4), Inches(5.3), Inches(1.3),
       color=RED, fill=RGBColor(0xF9,0xE4,0xE0), wdt=3)
txt(s, Inches(4.0), Inches(4.5), Inches(5.3), Inches(1.1), "COGNITIVE DEBT\n(the reps you skipped)",
    22, RED, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.9), Inches(6.0), Inches(11.5), Inches(1.2),
    "Let AI think FOR you instead of WITH you and the denominator grows —\n"
    "your real value shrinks, no matter how good the tool looks.", 19, INK, True, PP_ALIGN.CENTER)

# ---------- Slide 8: the tool ----------
s = slide()
txt(s, Inches(0.9), Inches(0.45), Inches(11.5), Inches(0.8),
    "🛠 A tool to try this week", 32, ACCENT, True)
txt(s, Inches(0.9), Inches(1.25), Inches(11.5), Inches(0.6),
    "The “AI-Free Rep” Protocol — 10 minutes, weekly", 23, INK, True)
steps = [
    "1.  Pick one core JUDGMENT task you now do with AI (a decision, diagnosis, plan).",
    "2.  Do it COLD first — write your own answer + reasoning before opening any AI. 5 min.",
    "3.  Then ask the AI. Compare. Where did you differ — and who was actually right?",
    "4.  Log the gap. If the AI changed your mind, was it correct, or did you just surrender?",
    "5.  Once a month go FULLY AI-free on one real task — your own private “Gartner assessment.”",
]
y = Inches(2.15)
for st in steps:
    doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), y, Inches(11.5), Inches(0.72),
           color=BLUE, fill=RGBColor(0xF0,0xF5,0xFB), wdt=2)
    txt(s, Inches(1.15), y+Inches(0.06), Inches(11.0), Inches(0.62), st, 16, INK, False,
        PP_ALIGN.LEFT, MSO_ANCHOR.MIDDLE)
    y = Emu(int(y + Inches(0.85)))
txt(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.8),
    "Bonus: prompt the AI “give 3 reasons my conclusion is wrong.” Sparring partner, not oracle.",
    16, ACCENT, True, PP_ALIGN.CENTER)

# ---------- Slide 9: closing / sources ----------
s = slide()
txt(s, Inches(0.9), Inches(0.7), Inches(11.5), Inches(1.0),
    "The takeaway", 40, ACCENT, True)
doodle(s, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.9), Inches(11.3), Inches(2.4),
       color=INK, fill=RGBColor(0xFF,0xF3,0xD9), wdt=3)
txt(s, Inches(1.4), Inches(2.1), Inches(10.5), Inches(2.0),
    "The winners won't use AI the most or the least —\n"
    "they'll use it DELIBERATELY, paying down cognitive debt\n"
    "with regular unaided reps so the judgment stays theirs.", 25, INK, True,
    PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
txt(s, Inches(0.9), Inches(4.55), Inches(11.5), Inches(2.2),
    "Sources: Gartner 2026 predictions · Anthropic deskilling RCT (via TNW) · Wharton “cognitive surrender” ·\n"
    "BCG cognitive-debt report · PwC 2026 AI Jobs Barometer · WEF Future of Jobs 2026 · ADP People at Work 2026 ·\n"
    "GMAC recruiter survey · Coursera micro-credentials 2026.  (Full links in the accompanying README.)",
    14, GREY, False, PP_ALIGN.CENTER)

OUT = "/home/user/my_professional_documents/weekly-briefings/2026-07-25-cognitive-debt-ai-deskilling/Cognitive_Debt_AI_Deskilling_2026.pptx"
prs.save(OUT)
print("saved deck with", len(prs.slides._sldIdLst), "slides ->", OUT)
